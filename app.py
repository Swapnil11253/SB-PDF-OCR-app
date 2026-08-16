import os
import glob
import re
import json
import copy
import io
import base64
import time
import shutil
import subprocess
import gradio as gr
import docx
try:
    import fitz  # PyMuPDF -- pure Python, no external binary, works on Windows
    FITZ_AVAILABLE = True
except ImportError:
    FITZ_AVAILABLE = False
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import parse_xml
from docx.oxml.ns import nsdecls, qn as docx_qn
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn as pptx_qn
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches as PPTX_Inches

# ======================================================================
# NATIVE LATEX -> OMML ENGINE  (latex2mathml + mathml2omml, pure Python)
# ======================================================================
# NOTE: Purana engine "pandoc" CLI par depend karta tha, lekin start.bat
# kabhi pandoc binary install hi nahi karta tha -> subprocess fail hota
# tha silently -> har equation LaTeX text ke fallback me chala jaata tha
# (isiliye Vector / Sum / Symbols convert nahi ho rahe the).
# Ab hum pure-Python pipeline use kar rahe: LaTeX -> MathML -> OMML.
# Isse koi external binary install karne ki zaroorat nahi.
try:
    import latex2mathml.converter as _l2m
    import mathml2omml as _m2o
    NATIVE_ENGINE_AVAILABLE = True
except ImportError:
    NATIVE_ENGINE_AVAILABLE = False

MATH_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"

# Track conversion stats for the QC report
_CONV_STATS = {"success": 0, "failed": 0, "failed_list": []}


def _reset_stats():
    _CONV_STATS["success"] = 0
    _CONV_STATS["failed"] = 0
    _CONV_STATS["failed_list"] = []


# OCR/formula-recognition (MinerU) very often misreads the dotless
# accented letter used under a hat/vec for unit-vector notation
# (i-hat/j-hat/k-hat, i.e. \hat{i} \hat{j} \hat{k}) as the DIGIT "1" --
# visually the thin dotless stroke under a circumflex/arrow looks a lot
# like "1" to the recognition model. A hat/vec/tilde/bar accent is never
# legitimately placed over a bare numeral in normal notation, so any
# \hat{1}, \vec{1}, etc. we see is almost certainly a misread unit
# vector. Standard vector notation always lists these in i, j, k order,
# so we restore the intended letters by cycling through that order
# across the equation.
_ACCENTED_ONE_RE = re.compile(
    r'\\(hat|widehat|vec|overrightarrow|tilde|widetilde)'
    r'(?:\s*\{\s*1\s*\}|\s+1(?![0-9]))'
)


def _fix_ocr_unit_vector_ones(s):
    letters = ['i', 'j', 'k']
    counter = {'n': 0}

    def _sub(m):
        cmd = m.group(1)
        letter = letters[counter['n'] % len(letters)]
        counter['n'] += 1
        return f'\\{cmd}{{{letter}}}'

    return _ACCENTED_ONE_RE.sub(_sub, s)


def sanitize_latex(latex_str):
    """
    Minimal, SAFE sanitizer. latex2mathml already natively understands
    \\left \\right \\Big \\vec \\overrightarrow \\otimes \\begin{array}/{matrix}
    etc, so we no longer strip that structure (the old code was stripping
    arrays/vectors/big-brackets, which is exactly why they broke).
    We only patch genuinely unsupported / legacy constructs.
    """
    s = latex_str.strip()

    # Legacy non-semantic font shells that sometimes trip parsers: {\rm xyz} / {\bf xyz}
    s = re.sub(r'\{\\rm\s+([^{}]+)\}', r'\\mathrm{\1}', s)
    s = re.sub(r'\{\\bf\s+([^{}]+)\}', r'\\mathbf{\1}', s)
    s = re.sub(r'\\hbox\s*\{([^{}]*)\}', r'\\text{\1}', s)

    # Common OCR/MinerU mis-escapes
    s = s.replace(r'\{cdot\}', r'\cdot').replace(r'\;', ' ').replace(r'\,', ' ')

    # Low-quality-scan OCR/formula-recognition models sometimes emit a bare
    # arrow character stuck to a variable (e.g. "L\u2192" or "L\u20d7")
    # instead of proper "\vec{L}" LaTeX. Normalize these into \vec{} so the
    # Accent-arrow pipeline below can pick them up correctly.
    s = re.sub(
        r'([A-Za-z])(?:_\{?([A-Za-z0-9]+)\}?)?[\u2192\u20d7]',
        lambda m: (r'\vec{' + m.group(1) + '_{' + m.group(2) + '}}') if m.group(2)
                  else (r'\vec{' + m.group(1) + '}'),
        s
    )

    # Repair OCR's i/j/k -> 1 misreads inside hat/vec unit-vector accents
    # (must run AFTER the bare-arrow normalization above, since that step
    # can itself produce fresh \vec{1} spans out of "1\u2192" OCR output).
    s = _fix_ocr_unit_vector_ones(s)

    # Collapse whitespace
    s = re.sub(r'\s+', ' ', s).strip()
    return s


def aggressive_sanitize_latex(s):
    """
    SECOND-PASS sanitizer, only tried if the first (safe) pass fails to
    convert. This is intentionally more invasive -- it trades a small risk
    of altering rare edge-case notation for a much higher chance of
    recovering a real OMML equation instead of the LaTeX-text fallback.
    Never called on the happy path, only as a retry.
    """
    s = re.sub(r'\\d?frac', r'\\frac', s)          # \dfrac / \tfrac -> \frac
    s = re.sub(r'\\varDelta', r'\\Delta', s)
    s = re.sub(r'\\varGamma', r'\\Gamma', s)
    s = re.sub(r'\\bold\{([^{}]*)\}', r'\\mathbf{\1}', s)
    s = re.sub(r'\\text\s*\{\s*\}', '', s)          # empty \text{}
    s = s.replace(r'\displaystyle', '').replace(r'\textstyle', '')
    s = s.replace(r'\nonumber', '').replace(r'\notag', '')
    s = re.sub(r'\\label\{[^{}]*\}', '', s)
    s = re.sub(r'\\tag\{[^{}]*\}', '', s)
    s = s.replace(r'\;', ' ').replace(r'\!', '').replace(r'\quad', '  ').replace(r'\qquad', '    ')
    # Balance an odd number of braces by stripping only genuinely stray
    # trailing/leading unmatched ones (never touch balanced pairs).
    if s.count('{') != s.count('}'):
        if s.count('{') > s.count('}'):
            s = s.rstrip('{')
        else:
            s = s.lstrip('}')
    return re.sub(r'\s+', ' ', s).strip()


def _repair_omml(omml_str):
    """
    Known bug in mathml2omml==0.0.2: for \\vec / \\overrightarrow (groupChr
    constructs) it emits an unclosed <m:groupChrPr>, closing it with
    </m:groupChr> instead of </m:groupChrPr>. This corrupts the XML and
    silently kills every vector-notation equation. Patch it here.
    """
    omml_str = re.sub(
        r'(<m:pos[^/]*/>)</m:groupChr>(<m:e>)',
        r'\1</m:groupChrPr>\2',
        omml_str
    )
    return omml_str


# Combining accent characters -> renders in Word's "Accent" gallery (tight,
# single-character vector arrow), matching what physics vector notation
# needs. Plain arrow characters (U+2192 etc) instead trigger Word's
# "Operator" stretchy-arrow rendering, which is wrong for \vec / \overrightarrow.
_ARROW_TO_ACCENT_CHAR = {
    '\u2192': '\u20d7',  # rightwards arrow  -> combining right arrow above (Vector)
    '\u2190': '\u20d6',  # leftwards arrow   -> combining left arrow above
    '\u2194': '\u20e1',  # left-right arrow  -> combining left-right arrow above
    '\u00af': '\u0305',  # macron (bar)      -> combining overline (\bar)
    '\u02c7': '\u030c',  # caron (check)     -> combining caron above (\check)
}

# mathml2omml renders \hat / \widehat / \tilde / \dot / \ddot / \breve as
# m:limUpp (Word's "Limit" gallery -- a wide bar/limit-style annotation
# sitting ABOVE the base on its own line), instead of the tight m:acc
# "Accent" gallery construct that Word's own Insert > Equation > Accent
# tool produces. This is why unit vectors like \hat{i} (i-hat / \hat{n})
# never looked like a native Word accent. The plain (non-combining)
# character mathml2omml puts inside <m:lim> tells us which accent it is;
# map that to the matching Unicode COMBINING accent character so the
# rewrite below can build a real m:acc.
_LIMUPP_LIM_CHAR_TO_ACCENT = {
    '^': '\u0302',      # circumflex / hat -> Word Accent gallery "Hat"
    '~': '\u0303',      # tilde
    '\u02d9': '\u0307',  # dot above
    '\u00a8': '\u0308',  # double dot (diaeresis) above -> \ddot
    '\u02d8': '\u0306',  # breve
}


def _vector_groupchr_to_accent(omath_el):
    """
    mathml2omml renders \\vec / \\overrightarrow as m:groupChr (Word's
    "Operator" stretchy-arrow category). Physics vector notation needs
    Word's "Accent" category instead (tight arrow directly over the
    variable) -- so we rewrite the relevant groupChr subtrees into m:acc.
    """
    for grp in omath_el.findall(f'.//{{{MATH_NS}}}groupChr'):
        pr = grp.find(f'{{{MATH_NS}}}groupChrPr')
        if pr is None:
            continue
        chr_el = pr.find(f'{{{MATH_NS}}}chr')
        pos_el = pr.find(f'{{{MATH_NS}}}pos')
        if chr_el is None or pos_el is None:
            continue
        chr_val = chr_el.get(f'{{{MATH_NS}}}val')
        pos_val = pos_el.get(f'{{{MATH_NS}}}val')
        if pos_val != 'top' or chr_val not in _ARROW_TO_ACCENT_CHAR:
            continue
        e_el = grp.find(f'{{{MATH_NS}}}e')
        if e_el is None:
            continue
        parent = grp.getparent()
        idx = list(parent).index(grp)
        acc = etree.Element(f'{{{MATH_NS}}}acc')
        acc_pr = etree.SubElement(acc, f'{{{MATH_NS}}}accPr')
        new_chr = etree.SubElement(acc_pr, f'{{{MATH_NS}}}chr')
        new_chr.set(f'{{{MATH_NS}}}val', _ARROW_TO_ACCENT_CHAR[chr_val])
        acc.append(e_el)
        parent.remove(grp)
        parent.insert(idx, acc)
    return omath_el


def _limupp_accent_to_acc(omath_el):
    """
    Rewrites m:limUpp subtrees produced for \\hat / \\widehat / \\tilde /
    \\dot / \\ddot / \\breve into m:acc, the SAME "Accent" gallery
    construct Word's own equation-accent tool uses -- so î / n̂ / etc.
    render as a tight accent sitting directly on the glyph, exactly like
    it would if built via Word's Accent tab, instead of the wider
    Limit-style spacing mathml2omml emits by default.
    """
    for limupp in omath_el.findall(f'.//{{{MATH_NS}}}limUpp'):
        lim_el = limupp.find(f'{{{MATH_NS}}}lim')
        e_el = limupp.find(f'{{{MATH_NS}}}e')
        if lim_el is None or e_el is None:
            continue
        t_el = lim_el.find(f'.//{{{MATH_NS}}}t')
        if t_el is None or t_el.text not in _LIMUPP_LIM_CHAR_TO_ACCENT:
            continue
        accent_char = _LIMUPP_LIM_CHAR_TO_ACCENT[t_el.text]
        parent = limupp.getparent()
        if parent is None:
            continue
        idx = list(parent).index(limupp)
        acc = etree.Element(f'{{{MATH_NS}}}acc')
        acc_pr = etree.SubElement(acc, f'{{{MATH_NS}}}accPr')
        new_chr = etree.SubElement(acc_pr, f'{{{MATH_NS}}}chr')
        new_chr.set(f'{{{MATH_NS}}}val', accent_char)
        acc.append(e_el)  # move the existing base element under the new acc
        parent.remove(limupp)
        parent.insert(idx, acc)
    return omath_el


def _fix_radical_missing_degree(omath_el):
    """
    mathml2omml omits <m:deg> entirely for plain \\sqrt{} (no nth-root
    index), instead of the required-by-schema <m:radPr><m:degHide/></m:radPr>
    + empty <m:deg/> pair. Word then renders an unwanted empty degree/index
    box next to ordinary square roots. Explicit \\sqrt[n]{} roots (which DO
    already have <m:deg> with real content) are left untouched.
    """
    for rad in omath_el.findall(f'.//{{{MATH_NS}}}rad'):
        if rad.find(f'{{{MATH_NS}}}deg') is not None:
            continue  # explicit nth-root -- already correct, don't touch
        rad_pr = etree.Element(f'{{{MATH_NS}}}radPr')
        deg_hide = etree.SubElement(rad_pr, f'{{{MATH_NS}}}degHide')
        deg_hide.set(f'{{{MATH_NS}}}val', '1')
        empty_deg = etree.Element(f'{{{MATH_NS}}}deg')
        rad.insert(0, empty_deg)
        rad.insert(0, rad_pr)
    return omath_el


def _convert_via_native_engine(cleaned_latex):
    if not NATIVE_ENGINE_AVAILABLE:
        return None
    mathml = _l2m.convert(cleaned_latex)
    omml_str = _repair_omml(_m2o.convert(mathml))
    wrapped = f'<root xmlns:m="{MATH_NS}">{omml_str}</root>'
    root = etree.fromstring(wrapped.encode("utf-8"))
    omath_el = root.find(f'.//{{{MATH_NS}}}oMath')
    if omath_el is not None:
        omath_el = _vector_groupchr_to_accent(omath_el)
        omath_el = _limupp_accent_to_acc(omath_el)
        omath_el = _fix_radical_missing_degree(omath_el)
    return omath_el


def _convert_via_pandoc(cleaned_latex):
    """Optional secondary fallback, used ONLY if pandoc happens to be
    installed on the system (fully optional, not required anymore)."""
    if not shutil.which("pandoc"):
        return None
    md_content = f"${cleaned_latex}$"
    try:
        process = subprocess.Popen(
            ["pandoc", "-f", "markdown+tex_math_dollars", "-t", "docx"],
            stdin=subprocess.PIPE, stdout=subprocess.PIPE, stderr=subprocess.PIPE
        )
        out, _ = process.communicate(input=md_content.encode("utf-8"), timeout=20)
        if out:
            import io, zipfile
            with zipfile.ZipFile(io.BytesIO(out)) as z:
                xml_content = z.read("word/document.xml")
                tree = etree.fromstring(xml_content)
                els = tree.xpath('//m:oMath | //m:oMathPara',
                                  namespaces={'m': MATH_NS})
                if els:
                    return els[0]
    except Exception:
        pass
    return None


def convert_latex_to_omml_native(latex_code):
    """
    Main entry point. Tries the native pure-Python engine first (fast,
    no external binary needed, handles vectors/sums/arrays/greek natively).
    Falls back to pandoc only if it's present on the system as a bonus.
    Returns an lxml <m:oMath> element, or None if every path fails.
    """
    cleaned = sanitize_latex(latex_code)
    if not cleaned:
        return None

    try:
        el = _convert_via_native_engine(cleaned)
        if el is not None:
            _CONV_STATS["success"] += 1
            return el
    except Exception:
        pass

    # RETRY with a more invasive sanitizer before giving up. Handles common
    # failure patterns (\dfrac, stray \label, unbalanced braces from OCR,
    # etc.) that the safe first pass deliberately leaves untouched.
    try:
        retry_latex = aggressive_sanitize_latex(cleaned)
        if retry_latex and retry_latex != cleaned:
            el = _convert_via_native_engine(retry_latex)
            if el is not None:
                _CONV_STATS["success"] += 1
                return el
    except Exception:
        pass

    el = _convert_via_pandoc(cleaned)
    if el is not None:
        _CONV_STATS["success"] += 1
        return el

    _CONV_STATS["failed"] += 1
    if len(_CONV_STATS["failed_list"]) < 15:
        _CONV_STATS["failed_list"].append(cleaned[:60])
    return None


# ======================================================================
# STRUCTURED-SPAN PROTECTION (fixes: multi-line vector/array equations
# and HTML tables getting silently CUT IN HALF by the naive \n\n / \n
# block-splitting that runs later -- the #1 cause of equations "staying
# in LaTeX" and tables/content being skipped even though the OMML engine
# itself works fine.)
# ======================================================================
_MATH_ENV_NAMES = (
    r"array|matrix|pmatrix|bmatrix|vmatrix|Vmatrix|align|aligned|"
    r"cases|gathered|split|equation|multline"
)
# ONE canonical pattern used both to protect spans before splitting AND
# to actually detect/convert math later, so the two can never drift apart.
MATH_SPAN_PATTERN = re.compile(
    r"(\$\$.*?\$\$"                      # $$ ... $$  (display math)
    r"|\\\[.*?\\\]"                       # \[ ... \]  (display math)
    r"|\\\(.*?\\\)"                       # \( ... \)  (inline math)
    r"|\$[^$\n]+?\$"                      # $ ... $    (inline math, single line)
    r"|\\begin\{(?:" + _MATH_ENV_NAMES + r")\}.*?\\end\{(?:" + _MATH_ENV_NAMES + r")\})",
    re.DOTALL
)
HTML_TABLE_PATTERN = re.compile(r"<table[\s\S]*?</table>", re.IGNORECASE)


def protect_structured_spans(text):
    """
    Run on the RAW markdown before any \\n\\n / \\n splitting happens.
    - HTML <table>...</table> blocks: collapse any blank lines inside them
      so the later \\n\\n block-split can't slice the table into pieces.
    - Math spans ($$..$$, \\[..\\], \\begin{env}..\\end{env} etc): flatten
      ALL internal newlines to spaces. LaTeX doesn't care about source
      formatting, and this guarantees the whole equation survives as one
      unbroken chunk through every later split.
    """
    def _collapse_blank_lines(m):
        return re.sub(r'\n\s*\n+', '\n', m.group(0))

    def _flatten(m):
        return re.sub(r'\s*\n\s*', ' ', m.group(0))

    text = HTML_TABLE_PATTERN.sub(_collapse_blank_lines, text)
    text = MATH_SPAN_PATTERN.sub(_flatten, text)
    return text


def _strip_latex_delimiters(span):
    """Strips the outer $$..$$ / \\[..\\] / \\(..\\) / $..$ wrapper off a
    matched math span (or is a no-op on already-bare LaTeX, e.g. a
    MinerU 'equation'-type block which has no delimiters to begin with)
    so every caller feeds convert_latex_to_omml_native() the same bare
    LaTeX regardless of which delimiter style the source used."""
    s = (span or "").strip()
    for opener, closer in (("$$", "$$"), (r"\[", r"\]"), (r"\(", r"\)")):
        if s.startswith(opener) and s.endswith(closer) and len(s) >= len(opener) + len(closer):
            return s[len(opener):-len(closer)].strip()
    if s.startswith("$") and s.endswith("$") and len(s) >= 2:
        return s[1:-1].strip()
    return s


# ======================================================================
# UNIVERSAL EDITED-WORD OBJECT TOKENS
# ======================================================================
OMML_TOKEN_RE = re.compile(r"@@OMML:([A-Za-z0-9+/=_-]+)@@")
IMAGE_TOKEN_RE = re.compile(r"@@IMG:([^@]+)@@")

def _encode_omml_token(omath_el):
    if omath_el is None:
        return ""
    return "@@OMML:" + base64.b64encode(etree.tostring(omath_el, encoding="utf-8")).decode("ascii") + "@@"

def _decode_omml_token(token):
    m = OMML_TOKEN_RE.fullmatch(token or "")
    if not m:
        return None
    try:
        return etree.fromstring(base64.b64decode(m.group(1)))
    except Exception:
        return None

def _omml_visible_text(omath_el):
    if omath_el is None:
        return ""
    return "".join((t.text or "") for t in omath_el.iter() if t.tag == "{%s}t" % MATH_NS).strip()

# Markdown-style emphasis markers -- this is the SAME convention MinerU
# emits inside content_list.json 'text' for source PDF bold/italic runs,
# and is also what _word_paragraph_runs_lossless() (edited-Word re-upload
# path) now wraps real w:b / w:i run formatting into -- so both PPTX
# source paths funnel through this ONE parser instead of the emphasis
# silently being dropped (literal "**" showing, or formatting vanishing).
_MD_STYLE_PATTERN = re.compile(
    r'\*\*\*(?P<bi>[^*]+?)\*\*\*'
    r'|\*\*(?P<b>[^*]+?)\*\*'
    r'|\*(?P<i>[^*]+?)\*'
)


def _split_markdown_style(text):
    """Splits plain text on **bold**/*italic*/***bold-italic*** markdown
    markers into an ordered list of (content, bold, italic) pieces. Text
    with no '*' at all (the overwhelming common case) is returned as a
    single untouched piece -- zero behavior change for plain content."""
    if not text:
        return []
    if '*' not in text:
        return [(text, False, False)]
    out = []
    pos = 0
    for m in _MD_STYLE_PATTERN.finditer(text):
        if m.start() > pos:
            out.append((text[pos:m.start()], False, False))
        if m.group('bi') is not None:
            out.append((m.group('bi'), True, True))
        elif m.group('b') is not None:
            out.append((m.group('b'), True, False))
        else:
            out.append((m.group('i'), False, True))
        pos = m.end()
    if pos < len(text):
        out.append((text[pos:], False, False))
    return out or [(text, False, False)]


def _iter_text_math_or_omml(raw_text):
    raw_text = raw_text or ""
    pieces = []
    pos = 0
    for m in OMML_TOKEN_RE.finditer(raw_text):
        if m.start() > pos:
            pieces.extend(_iter_text_and_math(raw_text[pos:m.start()]))
        el = _decode_omml_token(m.group(0))
        pieces.append(("omml", el) if el is not None else ("text", m.group(0), False, False))
        pos = m.end()
    if pos < len(raw_text):
        pieces.extend(_iter_text_and_math(raw_text[pos:]))
    return pieces or [("text", "", False, False)]

def _iter_text_and_math(raw_text):
    """Splits raw_text (plain text that may contain $..$, $$..$$, \\(..\\),
    \\[..\\], or \\begin{env}..\\end{env} math spans -- the SAME
    MATH_SPAN_PATTERN already used for detection elsewhere, so this can
    never drift out of sync with the rest of the pipeline) into an
    ordered list of ('text', str, bold, italic) / ('math', bare_latex,
    False, False) tuples. Every non-math 'text' piece is additionally run
    through _split_markdown_style() so **bold** / *italic* source
    formatting survives as real bold/italic runs instead of being
    collapsed to plain text. This is what lets every downstream consumer
    turn EVERY equation into a real, editable OMML equation instead of
    collapsing mixed text+math lines to Unicode-approximated plain text."""
    raw_text = raw_text or ""
    out = []
    pos = 0
    for m in MATH_SPAN_PATTERN.finditer(raw_text):
        if m.start() > pos:
            for chunk, bold, italic in _split_markdown_style(raw_text[pos:m.start()]):
                out.append(("text", chunk, bold, italic))
        out.append(("math", _strip_latex_delimiters(m.group(0)), False, False))
        pos = m.end()
    if pos < len(raw_text):
        for chunk, bold, italic in _split_markdown_style(raw_text[pos:]):
            out.append(("text", chunk, bold, italic))
    if not out:
        out.append(("text", "", False, False))
    return out


def parse_html_table_to_rows(html_str):
    """Parses a <table>...</table> HTML snippet (as emitted by MinerU's
    table-recognition model for complex/merged-cell tables) into a plain
    list-of-lists, same shape as the markdown pipe-table parser produces."""
    try:
        parser = etree.HTMLParser()
        tree = etree.fromstring(html_str, parser)
        rows = []
        for tr in tree.iter('tr'):
            cells = []
            for cell in tr:
                if cell.tag in ('td', 'th'):
                    text = ' '.join(''.join(cell.itertext()).split())
                    colspan = int(cell.get('colspan', 1) or 1)
                    cells.append(text)
                    for _ in range(colspan - 1):
                        cells.append("")  # keep column alignment for merged cells
            if cells:
                rows.append(cells)
        return rows
    except Exception:
        return []


# ======================================================================
# MCQ FORMATTING & STRICT SOURCE-NUMBERING PRESERVATION
# ======================================================================
# FIX: purana opt_pattern har digit 1-4 ko option-marker maan leta tha,
# isliye "12.", "21.", "3.4" jaise question numbers / decimals bhi
# galti se option jaisa split ho jaate the. Ab digit-markers ke liye
# strict lookaround lagaya hai (na pehle digit, na baad me digit) taaki
# sirf genuine standalone "1." / "2)" jaise option-labels hi pakde jaayein.
#
# FIX 2 (edited-Word reupload bug): jab user Word file me option-labels
# (A./B./C./D.) ko BOLD kar deta hai, to _word_paragraph_runs_lossless
# unhe markdown wrapper me lapet deta hai -- e.g. "A.  " ban jaata hai
# "**A.  **". Pehle OPT_PATTERN/QNUM_PATTERN sirf marker ke EXACT start
# (position 0) par match karte the, is liye "**A." me leading "**" ki
# wajah se match FAIL ho jaata tha -> poori MCQ "Subjective/Numerical"
# (qtype="other") bankar options TextBox 22-25 ki jagah sirf TextBox 21
# me chali jaati thi. Ab dono pattern ke shuru (aur end) me optional
# "\*{0,3}" allow kiya hai taaki bold/italic markdown wrapper ke bawajood
# asli marker sahi se detect ho, aur wrapper cleanly consume ho jaaye
# (koi orphan "**" bacha na rahe jo option text me leak ho jaaye).
OPT_PATTERN = re.compile(
    r"(?:^|(?<=\s))"
    r"\*{0,3}"
    r"(\([A-Da-d]\)|\[[A-Da-d]\]|[A-Da-d][\.\)]|(?<!\d)[1-4][\.\)](?!\d))"
    r"(?=\s|$)\s*\*{0,3}"
)
# FIX 3 (Swapnil): pehle sirf "1." / "Q1." / "1)" jaise EXACT
# "digit + period-ya-bracket" format hi question-number maana jaata tha --
# koi bhi doosra source-format ("Ques 1", "Q.No. 5", "(1)", "1:", "1-",
# "1]", ya sirf "Q1" bina kisi punctuation ke) is se match hi nahi karta
# tha, isliye wo questions kabhi MCQ template me nahi jaate the (silently
# "Subjective/Numerical" ya plain theory text bankar reh jaate the).
# Ab prefix bahut zyada flexible hai (Q/Ques/Que/Question/Qtn/Qs, "No."
# ke saath ya bina) aur closer punctuation bhi (. ) ] : -) sab allowed
# hai, PLUS "(1)" jaisa fully-parenthesized number bhi. Jab koi Q-prefix
# maujood hai to closer punctuation ki zaroorat nahi ("Q1 Find the..."
# bhi match karega) -- lekin PREFIX na ho to closer punctuation ab bhi
# COMPULSORY hai (jaise pehle), taaki plain numbers jo kisi paragraph ke
# beech me bina kisi marker ke aa jaate hain (e.g. "5 kg block...") galti
# se question-number na ban jaayein.
QNUM_PATTERN = re.compile(
    r"^\s*\*{0,3}"
    r"(Q(?:uestion|uestn|tn|ues|ue|s)?\.?\s*(?:No\.?)?\s*)?"
    r"\(?\s*(\d{1,3})\s*\)?"
    r"(?(1)(?:[\.\)\]:\-–]|(?=\s|$))|[\.\)\]:\-–])"
    r"\s*\*{0,3}(?!\d)"
)


def process_mcq_formatting(text):
    lines = text.split("\n")
    cleaned_blocks = []

    for line in lines:
        s = line.strip()
        if not s:
            cleaned_blocks.append("")
            continue

        # Strip forced auto-bullets (- or *) that MinerU sometimes injects
        s = re.sub(r"^[\-\*\+]\s+", "", s)

        # Force Horizontal Options to Vertical Line-by-Line
        matches = list(OPT_PATTERN.finditer(s))
        if len(matches) >= 2:
            vertical_options = []
            # CRITICAL: whatever text sits BEFORE the first marker is the
            # question stem (or a continuation of it) -- it must never be
            # dropped, even though it isn't itself "one of the matches".
            first_start = matches[0].start()
            if first_start > 0:
                prefix = s[:first_start].strip()
                if prefix:
                    vertical_options.append(prefix)
            for i in range(len(matches)):
                start = matches[i].start()
                end = matches[i + 1].start() if i + 1 < len(matches) else len(s)
                vertical_options.append(s[start:end].strip())
            cleaned_blocks.append("\n".join(vertical_options))
        else:
            cleaned_blocks.append(s)

    return "\n".join(cleaned_blocks)


def split_leading_marker(line):
    """
    Returns (marker, rest) if the line starts with a source question
    number (e.g. '12.', 'Q3)') or an option label (e.g. '(A)', 'B.').
    marker is None if no such pattern is found. We NEVER invent a number
    -- we only detect what's already in the source text, so the doc
    keeps exactly the numbering that was in the PDF.
    """
    m = QNUM_PATTERN.match(line)
    if m:
        return line[:m.end()].strip(), line[m.end():].strip()
    m2 = OPT_PATTERN.match(line)
    if m2 and m2.start() == 0:
        return line[:m2.end()].strip(), line[m2.end():].strip()
    return None, line


def _normalize_docx_marker(marker):
    """MS Word output must STRICTLY use 'N.' for question numbers (e.g.
    'Q3)', '12)', 'Q. No. 5' all become '3.', '12.', '5.') and 'X.' for
    option labels (e.g. '(A)', '[b]', 'C)' all become 'A.', 'B.', 'C.') --
    a single uppercase letter/number plus one period, nothing else. Only
    the DISPLAY marker's punctuation/casing is normalized here; the
    underlying digit or letter identity always comes from the source
    marker itself, never invented or renumbered."""
    if not marker:
        return marker
    qm = QNUM_PATTERN.match(marker)
    if qm:
        return f"{qm.group(2)}."
    om = OPT_PATTERN.match(marker)
    if om:
        label = re.sub(r"[^A-Za-z0-9]", "", om.group(1))
        if label.isdigit():
            idx = int(label) - 1
            letter = "ABCD"[idx] if 0 <= idx < 4 else label
        else:
            letter = label.upper()
        return f"{letter}."
    return marker


# ======================================================================
# ADVANCED WORD DOCUMENT BUILDER ENGINE
# ======================================================================
def _add_table_rows(doc, table_data):
    if not table_data:
        return
    n_cols = max(len(r) for r in table_data)
    table = doc.add_table(rows=len(table_data), cols=n_cols)
    table.style = 'Table Grid'
    for r_idx, row in enumerate(table_data):
        for c_idx in range(n_cols):
            val = row[c_idx] if c_idx < len(row) else ""
            cell = table.rows[r_idx].cells[c_idx]
            cell.text = val
            # 'Table Grid' style can carry its own font instead of
            # inheriting 'Normal' -- stamp Times New Roman 11pt on the
            # actual runs explicitly so table content matches the rest
            # of the document.
            for p in cell.paragraphs:
                for run in p.runs:
                    run.font.name = 'Times New Roman'
                    run.font.size = Pt(11)


def _add_image(doc, full_img_path):
    if not full_img_path or not os.path.exists(full_img_path):
        return
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(6)
    p.paragraph_format.space_after = Pt(6)
    run = p.add_run()
    try:
        run.add_picture(full_img_path, width=Inches(4.5))
    except Exception:
        pass


def _add_text_paragraph(doc, raw_line, bold_whole=False, font_size=11):
    """Renders one line of content as a Word paragraph: detects + bolds the
    SOURCE question number / option label (never invents numbering, and
    disables Word's own auto-numbering), then splits out and converts any
    inline/display math within the line to native OMML."""
    raw_line = raw_line.strip()
    if not raw_line:
        return

    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(4)
    p.paragraph_format.space_before = Pt(2)
    p.paragraph_format.line_spacing = 1.15
    # FIX: a paragraph whose ONLY content is a native OMML equation (e.g.
    # an MCQ option that's just "$x^2+1$") gets auto-centered by Word --
    # standalone m:oMath with no m:oMathPara/m:jc falls back to Word's
    # display-equation default, which is CENTER, not the document's
    # normal left alignment. Force LEFT explicitly so options/lines
    # stay flush with the rest of the question instead of drifting to
    # the middle of the page.
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT

    # Stop MS Word from creating auto-bullet/auto-number lists
    p_pr = p._p.get_or_add_pPr()
    p_pr.append(parse_xml(f'<w:numPr {nsdecls("w")}><w:numId w:val="0"/></w:numPr>'))

    marker, remainder = split_leading_marker(raw_line)
    if marker:
        marker = _normalize_docx_marker(marker)
        mrun = p.add_run(marker + "  ")
        mrun.font.bold = True
        mrun.font.name = 'Times New Roman'
        mrun.font.size = Pt(font_size)
        content_line = remainder
    else:
        content_line = raw_line

    parts = MATH_SPAN_PATTERN.split(content_line)
    for part in parts:
        if not part:
            continue

        stripped_part = part.strip()
        is_math = (
            (stripped_part.startswith("$") and stripped_part.endswith("$"))
            or (stripped_part.startswith(r"\[") and stripped_part.endswith(r"\]"))
            or (stripped_part.startswith(r"\(") and stripped_part.endswith(r"\)"))
            or stripped_part.startswith(r"\begin")
        )
        if is_math:
            clean_math = stripped_part
            for opener, closer in ((r"\[", r"\]"), (r"\(", r"\)")):
                if clean_math.startswith(opener) and clean_math.endswith(closer):
                    clean_math = clean_math[len(opener):-len(closer)]
            clean_math = clean_math.strip("$").strip()
            omml_element = convert_latex_to_omml_native(clean_math)

            if omml_element is not None:
                p._p.append(omml_element)
            else:
                # Robust fallback -- equation is NEVER dropped, worst
                # case it shows as raw LaTeX text so nothing is lost.
                run = p.add_run(f" {clean_math} ")
                run.font.italic = True
                run.font.color.rgb = RGBColor(0, 51, 102)
                run.font.name = 'Cambria Math'
        else:
            run = p.add_run(part)
            run.font.name = 'Times New Roman'
            run.font.size = Pt(font_size)
            run.font.bold = bold_whole


def _new_doc():
    doc = docx.Document()
    for section in doc.sections:
        section.top_margin = Inches(0.7)
        section.bottom_margin = Inches(0.7)
        section.left_margin = Inches(0.7)
        section.right_margin = Inches(0.7)
    # Document-wide default: Times New Roman, 11pt. This sets the base
    # 'Normal' style so anything that doesn't explicitly stamp its own
    # run font (e.g. table cell text from _add_table_rows) still comes
    # out Times New Roman 11pt rather than Word's own Calibri default.
    normal = doc.styles['Normal']
    normal.font.name = 'Times New Roman'
    normal.font.size = Pt(11)
    # East-Asian/complex-script font slots must be set too, or Word can
    # silently fall back to a different font for some glyphs even when
    # the main w:rFonts/@w:ascii says Times New Roman.
    rpr = normal.element.get_or_add_rPr()
    rfonts = rpr.find(docx_qn('w:rFonts'))
    if rfonts is None:
        rfonts = etree.SubElement(rpr, docx_qn('w:rFonts'))
    rfonts.set(docx_qn('w:ascii'), 'Times New Roman')
    rfonts.set(docx_qn('w:hAnsi'), 'Times New Roman')
    rfonts.set(docx_qn('w:cs'), 'Times New Roman')
    return doc


# ======================================================================
# STRUCTURED BUILDER (content_list.json) -- PREFERRED PATH
# ======================================================================
# MinerU's content_list.json gives every block a bbox [x0,y0,x1,y1]
# normalized to a 0-1000 page range, and pre-classified content (equation
# blocks carry clean LaTeX in "text" with no markdown-delimiter guessing
# needed; tables carry ready HTML in "table_body"). This sidesteps every
# markdown-regex fragility AND lets us enforce true reading order
# ourselves: complete the LEFT column top-to-bottom, THEN the RIGHT
# column top-to-bottom -- instead of trusting whatever order MinerU's
# own algorithm produced.
PAGE_FURNITURE_TYPES = set()  # STRICT: do not discard any PDF content
FULL_WIDTH_RATIO = 0.65  # bbox wider than this fraction of the PAGE'S OWN content width -> spans all columns

def _order_page_items_by_column(items):
    """STRICT PDF READING ORDER, page-by-page.

    Rules:
      1) A genuinely SINGLE-COLUMN page is read top-to-bottom.
      2) A genuinely TWO-COLUMN page is read:
            LEFT column top -> bottom
            THEN RIGHT column top -> bottom
         and only then the next PDF page is processed.
      3) We do NOT trust PDF/Word-style section-break, column-break or page-break
         metadata. The PDF page + bbox geometry is the source of truth.
      4) Full-width items (heading/banner/table/image spanning both columns) are
         kept in their vertical position: before the columns if above them, or
         after them if below them.
      5) Items without a bbox are NEVER discarded. They are appended in source
         order after positioned items.
    """
    parsed = []
    for source_index, it in enumerate(items):
        bbox = it.get("bbox")
        valid_bbox = None
        if isinstance(bbox, (list, tuple)) and len(bbox) == 4:
            try:
                x0, y0, x1, y1 = map(float, bbox)
                if x1 > x0 and y1 >= y0:
                    valid_bbox = [x0, y0, x1, y1]
            except (TypeError, ValueError):
                pass
        it["_bbox_f"] = valid_bbox
        it["_source_index"] = source_index
        parsed.append(it)

    positioned = [it for it in parsed if it["_bbox_f"] is not None]
    unpositioned = [it for it in parsed if it["_bbox_f"] is None]
    unpositioned.sort(key=lambda it: it["_source_index"])

    if not positioned:
        return unpositioned

    min_x = min(it["_bbox_f"][0] for it in positioned)
    max_x = max(it["_bbox_f"][2] for it in positioned)
    page_content_width = max(max_x - min_x, 1.0)

    # Very wide blocks are not useful for detecting columns. They are structural
    # full-width content (heading, wide table/image, etc.).
    body_candidates = []
    full_width = []
    for it in positioned:
        x0, y0, x1, y1 = it["_bbox_f"]
        width = x1 - x0
        if width >= 0.72 * page_content_width:
            full_width.append(it)
        else:
            body_candidates.append(it)

    def _sort_y(items_):
        return sorted(items_, key=lambda it: (
            it["_bbox_f"][1], it["_bbox_f"][0], it["_source_index"]
        ))

    # Detect a REAL two-column layout from a horizontal gap between item centers.
    # This is deliberately stricter than simply grouping overlapping x-ranges:
    # equations, fractions, bullets and small OCR fragments often create many
    # artificial x-bands, which previously caused the wrong reading order.
    two_column = False
    left_col, right_col = [], []

    if len(body_candidates) >= 4:
        # IMPORTANT: always provide an explicit key here.  Without it,
        # Python compares the second tuple element when two center-x values
        # are equal; that second element is a dict, which raises:
        # TypeError: '<' not supported between instances of 'dict' and 'dict'
        centers = sorted(
            (((it["_bbox_f"][0] + it["_bbox_f"][2]) / 2.0, it)
             for it in body_candidates),
            key=lambda pair: (pair[0], pair[1]["_source_index"])
        )

        # Find the largest gap between consecutive horizontal centers.
        best_gap = -1.0
        best_idx = -1
        for i in range(len(centers) - 1):
            gap = centers[i + 1][0] - centers[i][0]
            if gap > best_gap:
                best_gap = gap
                best_idx = i

        if best_idx >= 0:
            left_seed = [it for _, it in centers[:best_idx + 1]]
            right_seed = [it for _, it in centers[best_idx + 1:]]

            if left_seed and right_seed:
                left_center_max = max(
                    (it["_bbox_f"][0] + it["_bbox_f"][2]) / 2.0
                    for it in left_seed
                )
                right_center_min = min(
                    (it["_bbox_f"][0] + it["_bbox_f"][2]) / 2.0
                    for it in right_seed
                )

                # The gap must be substantial relative to page width.
                # 8% is intentionally conservative: a single-column page with
                # a heading/list indentation should not become "two columns".
                gap_ratio = best_gap / page_content_width

                def _vertical_span(group):
                    return (
                        min(it["_bbox_f"][1] for it in group),
                        max(it["_bbox_f"][3] for it in group)
                    )

                l_y0, l_y1 = _vertical_span(left_seed)
                r_y0, r_y1 = _vertical_span(right_seed)
                overlap = max(0.0, min(l_y1, r_y1) - max(l_y0, r_y0))
                shorter = max(1.0, min(l_y1 - l_y0, r_y1 - r_y0))
                vertical_overlap_ratio = overlap / shorter

                # Require both sides to contain a meaningful amount of content
                # and to occupy overlapping vertical regions. This prevents a
                # single-column page with a side image/footnote from being
                # misclassified as two-column.
                enough_items = min(len(left_seed), len(right_seed)) >= 2
                separated = (
                    gap_ratio >= 0.08
                    and right_center_min > left_center_max
                )
                two_column = (
                    enough_items
                    and separated
                    and vertical_overlap_ratio >= 0.25
                )

                if two_column:
                    # Use the geometric midpoint between the two center clusters
                    # as the column divider. This keeps every block, including
                    # equations/images, in exactly one column.
                    divider = (left_center_max + right_center_min) / 2.0
                    for it in body_candidates:
                        cx = (it["_bbox_f"][0] + it["_bbox_f"][2]) / 2.0
                        if cx < divider:
                            left_col.append(it)
                        else:
                            right_col.append(it)

                    # Safety: if the divider somehow creates an empty side,
                    # abandon two-column mode and use normal single-column order.
                    if not left_col or not right_col:
                        two_column = False
                        left_col, right_col = [], []

    if not two_column:
        # SINGLE-COLUMN MODE: nothing special is assumed. Everything is read
        # exactly by physical page position from top to bottom.
        all_positioned = _sort_y(positioned)
        return _merge_unpositioned_by_source_order(all_positioned, unpositioned)

    # TWO-COLUMN MODE:
    # Keep full-width content that physically sits above the column body first.
    # The main body is then STRICTLY left-to-bottom, right-to-bottom.
    # Full-width content below the body is appended after both columns.
    left_col = _sort_y(left_col)
    right_col = _sort_y(right_col)

    body_top = min(
        it["_bbox_f"][1] for it in (left_col + right_col)
    )
    body_bottom = max(
        it["_bbox_f"][3] for it in (left_col + right_col)
    )

    top_full = []
    middle_full = []
    bottom_full = []

    for it in full_width:
        y0, y1 = it["_bbox_f"][1], it["_bbox_f"][3]
        if y1 <= body_top:
            top_full.append(it)
        elif y0 >= body_bottom:
            bottom_full.append(it)
        else:
            middle_full.append(it)

    top_full = _sort_y(top_full)
    bottom_full = _sort_y(bottom_full)
    middle_full = _sort_y(middle_full)

    # A middle full-width block (for example a wide table between two column
    # sections) cannot safely be inserted into BOTH columns. Put it after the
    # left/right body only if it is actually below the column start; otherwise
    # keep it in source/vertical order at the end of the page. Most PDFs have
    # such blocks at the top or bottom, so this is the conservative choice.
    ordered = top_full + left_col + right_col + middle_full + bottom_full
    return _merge_unpositioned_by_source_order(ordered, unpositioned)


def _merge_unpositioned_by_source_order(ordered_positioned, unpositioned):
    """Items with no usable bbox (couldn't be geometrically placed) must
    still respect the ORIGINAL PDF source sequence relative to their
    neighbours -- previously they were all dumped at the very end of the
    page regardless of where they actually occurred, which is exactly the
    kind of out-of-sequence content the strict-reading-order requirement
    exists to prevent. Instead, splice each unpositioned item back in
    immediately after the positioned item that came right before it in
    the original source order (or at the very start of the page if
    nothing preceded it)."""
    if not unpositioned:
        return ordered_positioned
    unpositioned = sorted(unpositioned, key=lambda it: it["_source_index"])
    result = list(ordered_positioned)
    # Insert from the highest source_index down so earlier insertions
    # don't shift the target position of later ones.
    for u in reversed(unpositioned):
        u_idx = u["_source_index"]
        insert_at = len(result)
        for i, it in enumerate(result):
            if it["_source_index"] > u_idx:
                insert_at = i
                break
        result.insert(insert_at, u)
    return result

def build_docx_from_content_list(content_list_path, base_dir, images_dir, output_docx_path):
    _reset_stats()
    doc = _new_doc()

    with open(content_list_path, "r", encoding="utf-8") as f:
        content_list = json.load(f)

    # Group by page, preserving original index as a stable tiebreaker
    pages = {}
    for i, item in enumerate(content_list):
        page_idx = item.get("page_idx", 0)
        pages.setdefault(page_idx, []).append(item)

    skipped_types = {}
    total_seen = 0
    crash_recovered = 0

    for page_idx in sorted(pages.keys()):
        ordered_items = _order_page_items_by_column(pages[page_idx])

        for item in ordered_items:
            item_type = item.get("type", "text")

            if item_type in PAGE_FURNITURE_TYPES:
                continue  # reserved; currently empty so no extracted content is skipped

            total_seen += 1
            try:
                if item_type in ("text", "title"):
                    text = item.get("text", "")
                    is_title = item.get("text_level", 0) and item.get("text_level", 0) >= 1
                    # CRITICAL for scanned/OCR PDFs: the OCR/text-recognition
                    # model very often returns question + all 4 options merged
                    # into ONE string with no line breaks between them (e.g.
                    # "... (A) x (B) y (C) z (D) w"). Without this call, options
                    # never got split onto their own vertical lines.
                    text = text if is_title else process_mcq_formatting(text)
                    for line in text.split("\n"):
                        _add_text_paragraph(doc, line, bold_whole=bool(is_title),
                                             font_size=13 if is_title else 11)

                elif item_type == "equation":
                    latex = (item.get("text") or "").strip()
                    latex = latex.strip("$").strip()
                    if latex.startswith(r"\["):
                        latex = latex[2:]
                    if latex.endswith(r"\]"):
                        latex = latex[:-2]
                    latex = latex.strip()
                    if latex:
                        p = doc.add_paragraph()
                        p.paragraph_format.space_after = Pt(4)
                        p.paragraph_format.space_before = Pt(2)
                        omml_element = convert_latex_to_omml_native(latex)
                        if omml_element is not None:
                            p._p.append(omml_element)
                        else:
                            run = p.add_run(f" {latex} ")
                            run.font.italic = True
                            run.font.color.rgb = RGBColor(0, 51, 102)
                            run.font.name = 'Cambria Math'
                    # equations may also carry a rendered snapshot image; skip it,
                    # the OMML/text fallback above already captures the content.

                elif item_type == "table":
                    table_html = item.get("table_body", "")
                    table_data = parse_html_table_to_rows(table_html) if table_html else []
                    if table_data:
                        _add_table_rows(doc, table_data)
                    else:
                        img_path = item.get("img_path")
                        if img_path and os.path.exists(os.path.join(base_dir, img_path)):
                            _add_image(doc, os.path.join(base_dir, img_path))
                        else:
                            # Nothing structured recoverable -- NEVER go silent.
                            # Emit whatever raw text exists, else a visible flag
                            # so a human reviewer knows content sat here.
                            raw_fallback = item.get("text") or table_html
                            if raw_fallback:
                                _add_text_paragraph(doc, raw_fallback)
                            else:
                                _add_text_paragraph(
                                    doc,
                                    "[TABLE DETECTED BUT NOT RECOVERABLE -- manual review needed, page %s]" % page_idx
                                )
                    for cap in item.get("table_caption", []) or []:
                        _add_text_paragraph(doc, cap)

                elif item_type in ("image", "chart"):
                    img_path = item.get("img_path")
                    if img_path and os.path.exists(os.path.join(base_dir, img_path)):
                        _add_image(doc, os.path.join(base_dir, img_path))
                    else:
                        _add_text_paragraph(
                            doc,
                            "[IMAGE/DIAGRAM DETECTED BUT FILE MISSING -- manual review needed, page %s]" % page_idx
                        )
                    for cap in item.get("image_caption", []) or item.get("chart_caption", []) or []:
                        _add_text_paragraph(doc, cap)

                elif item_type in ("code", "list"):
                    # Render as plain text so content is never silently dropped
                    text = item.get("code_body") or "\n".join(item.get("list_items", []))
                    for line in (text or "").split("\n"):
                        _add_text_paragraph(doc, line)

                else:
                    skipped_types[item_type] = skipped_types.get(item_type, 0) + 1
                    # Unknown type -- still emit whatever text exists so nothing
                    # is silently lost.
                    fallback_text = item.get("text") or ""
                    if fallback_text:
                        _add_text_paragraph(doc, fallback_text)
                    else:
                        _add_text_paragraph(
                            doc,
                            "[UNRECOGNIZED CONTENT TYPE '%s' WITH NO TEXT -- manual review needed, page %s]"
                            % (item_type, page_idx)
                        )

            except Exception as e:
                # A single malformed item (bad bbox, corrupt LaTeX, broken
                # image file, etc.) must NEVER take down the whole document.
                # Recover by dumping whatever raw text we have as plain text
                # and moving on -- so the rest of the document still builds.
                crash_recovered += 1
                try:
                    raw = item.get("text") or item.get("table_body") or ""
                    note = "[RECOVERED AFTER ERROR (%s) -- manual review needed, page %s] %s" % (
                        type(e).__name__, page_idx, raw
                    )
                    _add_text_paragraph(doc, note)
                except Exception:
                    pass  # even the recovery paragraph failed -- give up on this one item only

    doc.save(output_docx_path)
    stats = dict(_CONV_STATS)
    stats["skipped_types"] = skipped_types
    stats["total_blocks_seen"] = total_seen
    stats["crash_recovered"] = crash_recovered
    return stats


# ======================================================================
# MARKDOWN-BASED BUILDER (fallback path, used only if content_list.json
# is not produced by the installed MinerU version)
# ======================================================================
def build_enterprise_docx(md_path, images_dir, output_docx_path):
    _reset_stats()
    doc = _new_doc()

    with open(md_path, "r", encoding="utf-8") as f:
        raw_md = f.read()

    # Protect multi-line equations & HTML tables BEFORE any \n\n / \n
    # splitting happens, so they can never be cut in half.
    raw_md = protect_structured_spans(raw_md)

    processed_text = process_mcq_formatting(raw_md)

    img_pattern = re.compile(r"!\[.*?\]\((.*?)\)")
    math_pattern = MATH_SPAN_PATTERN
    table_row_pattern = re.compile(r"^\s*\|.*\|\s*$")

    blocks = processed_text.split("\n\n")

    for block in blocks:
        block_str = block.strip()
        if not block_str:
            continue

        # 1. Images & Diagrams
        img_match = img_pattern.search(block_str)
        if img_match:
            img_rel_path = img_match.group(1)
            full_img_path = os.path.join(os.path.dirname(md_path), img_rel_path)
            if not os.path.exists(full_img_path):
                full_img_path = os.path.join(images_dir, os.path.basename(img_rel_path))
            _add_image(doc, full_img_path)
            continue

        # 2. Tables (HTML <table> first -- MinerU's table-recognition model
        #    emits HTML for complex/merged-cell tables, which the old code
        #    never handled at all, so those tables were silently rendered
        #    as broken raw-tag text or skipped-looking content.)
        if "<table" in block_str.lower():
            table_data = parse_html_table_to_rows(block_str)
            if table_data:
                _add_table_rows(doc, table_data)
                continue

        # 2b. Tables (Markdown pipe-table syntax)
        lines = block_str.split("\n")
        if all(table_row_pattern.match(l) for l in lines if l.strip()):
            table_data = []
            for line in lines:
                if "---" in line:
                    continue
                cells = [c.strip() for c in line.split("|")[1:-1]]
                if cells:
                    table_data.append(cells)
            if table_data:
                _add_table_rows(doc, table_data)
                continue

        # 3. Text / MCQs / Math -> possibly multiple lines per block
        for raw_line in block_str.split("\n"):
            _add_text_paragraph(doc, raw_line)

    doc.save(output_docx_path)
    return dict(_CONV_STATS)


# ======================================================================
# CORRUPT SYMBOL-FONT DETECTION
# ======================================================================
# Coaching-institute solution PDFs (ALLEN/Kota style, MathType/Equation
# Editor exports) very often embed equation glyphs through fonts named
# "SymbolMT" / "MT-Extra" / "Wingdings" etc. with Identity-H CID encoding
# whose ToUnicode CMap is broken. Text extracted from these fonts comes
# out as WRONG characters (e.g. "x" -> "A-umlaut", "Sum" -> "a-ring") and
# in the wrong reading order -- no amount of "text extraction" fixing will
# help here, because the source data itself is unrecoverable as text.
# The only reliable fix is to bypass the text layer entirely and run the
# page through the full visual OCR + formula-recognition pipeline. We
# auto-detect this risk up front and force that pipeline on.
RISKY_FONT_KEYWORDS = (
    # Only fonts with a KNOWN history of broken/missing ToUnicode cmaps
    # (old MS-Office "Insert Equation" era fonts). Deliberately excludes
    # LaTeX's own Computer Modern math fonts (cmsy/cmex/cmmi) -- those
    # extract cleanly in the vast majority of real PDFs and forcing OCR
    # on them was a pure speed/accuracy regression: full-page OCR is both
    # SLOWER and LESS accurate than reading an intact embedded text layer.
    "symbol", "mt-extra", "mtextra", "wingding",
    "esstix", "euclid", "mathtype", "msam", "msbm", "rsfs", "blondieburton"
)


def detect_symbol_font_risk(pdf_path):
    """Returns (risk_detected: bool, matched_fonts: list[str])."""
    if not FITZ_AVAILABLE:
        return False, []
    matched = set()
    try:
        doc = fitz.open(pdf_path)
        for page in doc:
            for f in page.get_fonts():
                basefont = (f[3] or "").lower()
                encoding = (f[5] or "").lower()
                if "identity-h" in encoding and any(k in basefont for k in RISKY_FONT_KEYWORDS):
                    matched.add(f[3])
        doc.close()
    except Exception:
        return False, []
    return (len(matched) > 0), sorted(matched)


# ======================================================================
# MAIN PROCESSOR PIPELINE
# ======================================================================
LANG_MAP = {
    "Auto / English (default)": None,
    "Hindi / Devanagari": "devanagari",
    "Chinese (best OCR base model)": "ch",
    "Tamil": "ta",
    "Telugu": "te",
}


EXTRACTION_MODES = {
    "Auto-detect (recommended)": "auto_detect",
    "Force Full OCR (garbled symbols / scanned / low quality)": "force_ocr",
    "Fast Text Mode (clean digital PDFs only)": "force_txt",
}


# ======================================================================
# WORD/PDF -> POWERPOINT (TEMPLATE-DRIVEN) ENGINE
# ======================================================================
# Design, driven by inspecting the user's actual "Format.pptx":
#   Slide 0 (Theory blueprint): 'Rectangle: Rounded Corners 1' = heading,
#            'TextBox 1' = body content, 'Group 2' = static logo (untouched)
#   Slide 1 (Q&A blueprint):    'TextBox 21' = question, 'TextBox 22..25'
#            = options A..D, 'Group 3/9/13/17' = static A/B/C/D circle
#            badges (untouched), 'TextBox 8' = static "QUESTION" label
#
# Both blueprint slides are NEVER written into directly -- every real
# slide in the output is a fresh XML-level duplicate of whichever
# blueprint is needed, in the exact order content appears in the source
# (theory and questions interleave arbitrarily), so the final deck's
# slide order always matches the source document. The two original
# blueprint slides are deleted from the output at the very end.

THEORY_HEADING_SHAPE = "Rectangle: Rounded Corners 1"
THEORY_BODY_SHAPE = "TextBox 1"
Q_STEM_SHAPE = "TextBox 21"
Q_OPTION_SHAPES = ["TextBox 22", "TextBox 23", "TextBox 24", "TextBox 25"]
Q_OPTION_GROUP_NAMES = ["Group 3", "Group 9", "Group 13", "Group 17"]

# --- LaTeX -> Unicode best-effort renderer -----------------------------
# PowerPoint text boxes (unlike Word) have no OMML math-run support via
# python-pptx, and this template defines single plain-text boxes (no
# dedicated equation-image placeholder). To guarantee NOTHING is ever
# silently dropped, every equation is converted to the closest readable
# Unicode approximation instead. Any LaTeX command with no mapping is
# kept visible as "[commandname]" rather than vanishing.
_GREEK_MAP = {
    "alpha": "α", "beta": "β", "gamma": "γ", "Gamma": "Γ", "delta": "δ",
    "Delta": "Δ", "epsilon": "ε", "varepsilon": "ε", "zeta": "ζ", "eta": "η",
    "theta": "θ", "Theta": "Θ", "vartheta": "ϑ", "iota": "ι", "kappa": "κ",
    "lambda": "λ", "Lambda": "Λ", "mu": "μ", "nu": "ν", "xi": "ξ", "Xi": "Ξ",
    "pi": "π", "Pi": "Π", "rho": "ρ", "sigma": "σ", "Sigma": "Σ", "tau": "τ",
    "upsilon": "υ", "phi": "φ", "varphi": "φ", "Phi": "Φ", "chi": "χ",
    "psi": "ψ", "Psi": "Ψ", "omega": "ω", "Omega": "Ω",
}
_SYMBOL_MAP = {
    "times": "×", "div": "÷", "pm": "±", "mp": "∓", "cdot": "·", "circ": "∘",
    "infty": "∞", "partial": "∂", "nabla": "∇", "sum": "Σ", "int": "∫",
    "oint": "∮", "approx": "≈", "neq": "≠", "leq": "≤", "geq": "≤",
    "propto": "∝", "rightarrow": "→", "to": "→", "leftarrow": "←",
    "Rightarrow": "⇒", "leftrightarrow": "↔", "cdots": "⋯", "ldots": "…",
    "degree": "°", "perp": "⊥", "parallel": "∥", "angle": "∠", "prime": "′",
    "sim": "∼", "equiv": "≡", "forall": "∀", "exists": "∃", "in": "∈",
    "hbar": "ℏ", "ohm": "Ω", "checkmark": "✓",
}
_SUPERSCRIPT_MAP = {
    "0": "⁰", "1": "¹", "2": "²", "3": "³", "4": "⁴", "5": "⁵", "6": "⁶",
    "7": "⁷", "8": "⁸", "9": "⁹", "+": "⁺", "-": "⁻", "=": "⁼", "(": "⁽",
    ")": "⁾", "n": "ⁿ", "i": "ⁱ", "o": "°",
}
_SUBSCRIPT_MAP = {
    "0": "₀", "1": "₁", "2": "₂", "3": "₃", "4": "₄", "5": "₅", "6": "₆",
    "7": "₇", "8": "₈", "9": "₉", "+": "₊", "-": "₋", "=": "₌", "(": "₍",
    ")": "₎", "a": "ₐ", "e": "ₑ", "h": "ₕ", "k": "ₖ", "l": "ₗ", "m": "ₘ",
    "n": "ₙ", "o": "ₒ", "p": "ₚ", "s": "ₛ", "t": "ₜ", "u": "ᵤ", "v": "ᵥ",
    "x": "ₓ", "r": "ᵣ",
}


def _supersub_convert(inner, mapping):
    if inner and all(c in mapping for c in inner):
        return "".join(mapping[c] for c in inner)
    return None


def _consume_brace_group(s, i):
    """s[i] must be '{'. Returns (content_without_braces, index_after_closing_brace),
    correctly handling nested braces (e.g. \\dfrac{Q_{enc}}{\\epsilon_0})."""
    if i >= len(s) or s[i] != '{':
        return None, i
    depth = 0
    j = i
    while j < len(s):
        if s[j] == '{':
            depth += 1
        elif s[j] == '}':
            depth -= 1
            if depth == 0:
                return s[i + 1:j], j + 1
        j += 1
    return s[i + 1:], len(s)  # unbalanced brace -- take the rest rather than crash


def _replace_frac_commands(s):
    out = []
    i = 0
    pattern = re.compile(r'\\d?frac')
    while True:
        m = pattern.search(s, i)
        if not m:
            out.append(s[i:])
            break
        out.append(s[i:m.start()])
        j = m.end()
        while j < len(s) and s[j] == ' ':
            j += 1
        if j < len(s) and s[j] == '{':
            num, j2 = _consume_brace_group(s, j)
            while j2 < len(s) and s[j2] == ' ':
                j2 += 1
            if j2 < len(s) and s[j2] == '{':
                den, j3 = _consume_brace_group(s, j2)
                out.append('(%s)/(%s)' % (num, den))
                i = j3
                continue
        out.append('[frac]')  # malformed args -- keep visible, never silently drop
        i = j
    return ''.join(out)


def _replace_single_arg_commands(s, command_names, transform):
    pattern = re.compile(r'\\(?:' + '|'.join(command_names) + r')\s*')
    out = []
    i = 0
    while True:
        m = pattern.search(s, i)
        if not m:
            out.append(s[i:])
            break
        out.append(s[i:m.start()])
        j = m.end()
        if j < len(s) and s[j] == '{':
            arg, j2 = _consume_brace_group(s, j)
            out.append(transform(arg))
            i = j2
        elif j < len(s):
            out.append(transform(s[j]))
            i = j + 1
        else:
            i = j
    return ''.join(out)


def _replace_supersub(s, marker, mapping, wrap_open, wrap_close):
    """Scans for ^ or _ and brace-matches the argument (so nested markup
    like r^{2n} or Q_{enc} works), converting to Unicode where every
    character in the argument is mappable, else a readable ^(...)/_(...) 
    fallback -- never silently dropped."""
    out = []
    i = 0
    while True:
        idx = s.find(marker, i)
        if idx == -1:
            out.append(s[i:])
            break
        out.append(s[i:idx])
        j = idx + 1
        if j < len(s) and s[j] == '{':
            arg, j2 = _consume_brace_group(s, j)
            conv = _supersub_convert(arg, mapping)
            out.append(conv if conv else wrap_open + arg + wrap_close)
            i = j2
        elif j < len(s):
            ch = s[j]
            conv = _supersub_convert(ch, mapping)
            out.append(conv if conv else wrap_open + ch + wrap_close)
            i = j + 1
        else:
            out.append(marker)
            i = j
    return ''.join(out)


def latex_to_display_text(latex):
    """Best-effort LaTeX -> plain Unicode conversion for PPTX text boxes.
    Never raises, never returns empty for non-empty input, never silently
    drops an unrecognized command (falls back to '[command]')."""
    if not latex:
        return ""
    s = latex.strip()
    s = re.sub(r'^\${1,2}|\${1,2}$', '', s)
    s = re.sub(r'^\\\(|\\\)$', '', s)
    s = re.sub(r'^\\\[|\\\]$', '', s)
    s = s.strip()

    s = _replace_frac_commands(s)
    s = _replace_single_arg_commands(s, ['sqrt'], lambda a: '√(' + a + ')')
    s = _replace_single_arg_commands(s, ['vec'], lambda a: a + '\u20d7')
    s = _replace_single_arg_commands(s, ['hat'], lambda a: a + '\u0302')
    s = _replace_single_arg_commands(s, ['bar'], lambda a: a + '\u0305')
    s = _replace_single_arg_commands(s, ['dot'], lambda a: a + '\u0307')
    s = _replace_single_arg_commands(s, ['text', 'mathbf', 'mathrm'], lambda a: a)
    s = re.sub(r'\\left|\\right', '', s)

    for name, ch in sorted(_GREEK_MAP.items(), key=lambda kv: -len(kv[0])):
        s = re.sub(r'\\' + name + r'(?![A-Za-z])', ch, s)
    for name, ch in sorted(_SYMBOL_MAP.items(), key=lambda kv: -len(kv[0])):
        s = re.sub(r'\\' + name + r'(?![A-Za-z])', ch, s)

    s = _replace_supersub(s, '^', _SUPERSCRIPT_MAP, '^(', ')')
    s = _replace_supersub(s, '_', _SUBSCRIPT_MAP, '_(', ')')

    # Anything left over with a backslash-command: keep it VISIBLE rather
    # than silently vanishing.
    s = re.sub(r'\\([a-zA-Z]+)', r'[\1]', s)
    s = s.replace('{', '').replace('}', '')
    return re.sub(r'\s+', ' ', s).strip()


def render_mixed_text(text):
    """Human-readable preview for LaTeX plus lossless native Word-OMML tokens."""
    if not text:
        return ""
    def _replace_omml(m):
        el = _decode_omml_token(m.group(0))
        return _omml_visible_text(el) if el is not None else m.group(0)
    text = OMML_TOKEN_RE.sub(_replace_omml, text)
    return MATH_SPAN_PATTERN.sub(lambda m: latex_to_display_text(m.group(0)), text)


# --- Slide-duplication (python-pptx has no native "duplicate slide") ---
def duplicate_slide(prs, index):
    """XML-level duplicate of prs.slides[index], appended at the end.
    Re-creates every relationship (critical for the logo picture) on the
    new slide part instead of copying raw rIds, which is what actually
    keeps the file from opening as 'corrupted / needs repair' in
    PowerPoint -- a raw XML deepcopy alone breaks embedded images."""
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    rid_map = {}
    for rel_id, rel in source.part.rels.items():
        if rel.is_external:
            new_rid = dest.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rid = dest.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rel_id] = new_rid

    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

    def remap_rids(el):
        for node in el.iter():
            for attr_name in ('embed', 'link', 'id'):
                val = node.get('{%s}%s' % (r_ns, attr_name))
                if val and val in rid_map:
                    node.set('{%s}%s' % (r_ns, attr_name), rid_map[val])

    for shp in source.shapes:
        el = copy.deepcopy(shp._element)
        remap_rids(el)
        dest.shapes._spTree.append(el)
    return dest


def _delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])


def _find_shape(container, name):
    for shp in container.shapes:
        if shp.name == name:
            return shp
    return None


def _remove_shape(shape):
    el = shape._element
    parent = el.getparent()
    if parent is not None:
        parent.remove(el)


def _set_shape_text(shape, text):
    """Replace a template textbox's content, cloning the EXACT original
    run/paragraph XML formatting (font, size, color, alignment) instead of
    python-pptx defaults, so output matches the template precisely. Extra
    lines become extra paragraphs using the same cloned formatting."""
    tf = shape.text_frame
    first_p = tf.paragraphs[0]
    pPr_src = first_p._p.find(pptx_qn('a:pPr'))
    rPr_src = None
    if first_p.runs:
        rPr_src = first_p.runs[0]._r.find(pptx_qn('a:rPr'))

    tf.clear()
    lines = (text or "").split("\n") or [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if pPr_src is not None:
            existing_pPr = p._p.find(pptx_qn('a:pPr'))
            if existing_pPr is not None:
                p._p.remove(existing_pPr)
            p._p.insert(0, copy.deepcopy(pPr_src))
        run = p.add_run()
        run.text = line
        if rPr_src is not None:
            existing_rPr = run._r.find(pptx_qn('a:rPr'))
            if existing_rPr is not None:
                run._r.remove(existing_rPr)
            run._r.insert(0, copy.deepcopy(rPr_src))


_HEADING_FONT_CANDIDATES_BOLD = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "DejaVuSans-Bold.ttf", "arialbd.ttf", "Arial Bold.ttf",
    "C:\\Windows\\Fonts\\arialbd.ttf", "C:\\Windows\\Fonts\\calibrib.ttf",
]
_HEADING_FONT_CANDIDATES_REGULAR = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "DejaVuSans.ttf", "arial.ttf", "Arial.ttf",
    "C:\\Windows\\Fonts\\arial.ttf", "C:\\Windows\\Fonts\\calibri.ttf",
]


def _measure_text_width_emu(text, font_size_pt, bold=True):
    """Best-effort text width estimate in EMU. Used to proactively resize
    a heading shape to fit its text IN THE SAVED FILE ITSELF, rather than
    relying only on PowerPoint's on-open autofit recalculation -- some
    viewers (and PowerPoint before the file is first touched/edited)
    don't always re-run that pass, which is exactly how a long heading
    used to visibly spill outside its shape. Uses real font metrics via
    PIL when a usable font file is found, else a calibrated average-
    glyph-width heuristic (never fails, always returns something)."""
    text = text or ""
    if not text:
        return 0
    try:
        from PIL import ImageFont
        px_size = max(8, int(round(font_size_pt * 96.0 / 72.0)))
        candidates = _HEADING_FONT_CANDIDATES_BOLD if bold else _HEADING_FONT_CANDIDATES_REGULAR
        for c in candidates:
            try:
                font = ImageFont.truetype(c, px_size)
                bbox = font.getbbox(text)
                width_px = bbox[2] - bbox[0]
                if width_px > 0:
                    return int((width_px / 96.0) * 914400)
            except Exception:
                continue
    except Exception:
        pass
    # Calibrated avg glyph-width fallback, ONLY used when no usable font
    # file is found for PIL to measure with. Bumped up from the previous
    # 0.58/0.52 factors -- those under-measured real heading fonts on
    # several machines (no matching TTF found -> fallback path taken),
    # which is exactly why headings were spilling outside their pill
    # shape: the shape was sized to a text-width estimate that was too
    # narrow. Deliberately erring wide here is safe (a heading pill that
    # is a little too big still looks fine); erring narrow is not (text
    # visibly clipped outside the shape).
    factor = 0.66 if bold else 0.60
    width_in = (font_size_pt / 72.0) * factor * len(text)
    return int(width_in * 914400)


def _force_heading_autofit(shape, text):
    """Explicitly forces BOTH halves of the requested autofit behavior:
    'Resize shape to fit text' ON (spAutoFit) and 'Wrap text in shape'
    OFF (wrap=none) -- set unconditionally every time (not just relied on
    from the template) because editing a shape's text via python-pptx can
    otherwise leave these in an inconsistent state. On top of the XML
    flags, it ALSO proactively grows the shape's actual saved width/height
    to fit `text` right now (with a small safety margin), so the heading
    never appears to spill outside its box even in a viewer that doesn't
    recompute autofit before the file is first opened/edited."""
    tf = shape.text_frame
    tf.word_wrap = False
    try:
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    except Exception:
        pass

    font_size_pt = _first_run_font_size_pt(shape)
    bold = False
    if tf.paragraphs and tf.paragraphs[0].runs:
        bold = bool(tf.paragraphs[0].runs[0].font.bold)

    lines = (text or "").split("\n") or [""]
    widest_emu = max((_measure_text_width_emu(l, font_size_pt, bold) for l in lines), default=0)

    bodyPr = tf._bodyPr
    l_ins = bodyPr.lIns if bodyPr.lIns is not None else 91440
    r_ins = bodyPr.rIns if bodyPr.rIns is not None else 91440
    t_ins = bodyPr.tIns if bodyPr.tIns is not None else 45720
    b_ins = bodyPr.bIns if bodyPr.bIns is not None else 45720

    SAFETY = 1.20  # generous cushion -- a bit of extra pill-padding beats any clipped heading
    new_width = int(widest_emu * SAFETY) + l_ins + r_ins

    line_h_emu = int((font_size_pt / 72.0) * 1.30 * 914400)
    new_height = line_h_emu * max(1, len(lines)) + t_ins + b_ins

    # Always SET (not just grow) -- "Resize shape to fit text" means the
    # box tracks the text both up AND down. A grow-only check left a
    # short heading sitting inside the template's original (often much
    # wider) box, which is exactly why autofit looked like it wasn't
    # applying at all.
    MIN_WIDTH = int(0.5 * 914400)
    MIN_HEIGHT = line_h_emu + t_ins + b_ins
    shape.width = max(new_width, MIN_WIDTH)
    shape.height = max(new_height, MIN_HEIGHT)

    # Never let the heading pill grow wider than the slide itself -- if
    # the estimate above would push it off the right edge, cap the width
    # to what's actually available from the shape's left position and
    # switch wrapping back ON so the (rare, very long) heading wraps to a
    # second line INSIDE the shape instead of running off the slide.
    try:
        slide_w = shape.part.package.presentation_part.presentation.slide_width
    except Exception:
        slide_w = None
    if slide_w:
        right_margin = int(0.15 * 914400)
        max_w = max(MIN_WIDTH, slide_w - shape.left - right_margin)
        if shape.width > max_w:
            shape.width = max_w
            tf.word_wrap = True
            # With wrapping back on, height must grow to fit however many
            # lines the now-capped width forces the text onto.
            est_lines = max(1, len(lines), int((widest_emu * SAFETY) / max(1, max_w - l_ins - r_ins)) + 1)
            shape.height = max(line_h_emu * est_lines + t_ins + b_ins, MIN_HEIGHT)


def _set_heading_text(shape, text):
    """Set heading text on the rounded-pill heading shape and let its
    WIDTH hug the new text -- grows for a long heading, shrinks back for
    a short one -- exactly like PowerPoint's own 'Resize shape to fit
    text'. LEFT anchor and HEIGHT stay pinned to the template's values
    (this is a fixed-height, single-line badge; only its horizontal
    extent should flow with however long the heading text is).

    Previously the shape's original width/height were force-restored
    AFTER the text was set (and auto_size was set to shrink the TEXT to
    fit the fixed shape instead of growing the shape), which is exactly
    why the pill never tracked the heading's length.
    """
    left, top, height = shape.left, shape.top, shape.height
    _set_shape_text_with_math(shape, text)
    shape.left, shape.top, shape.height = left, top, height

    # Reuses the same width/height auto-fit math already relied on for
    # the theory-slide heading elsewhere, then pins TOP back to the
    # template's fixed value so only width (and, for a very long heading
    # that had to wrap onto a 2nd line to stay inside the slide, height)
    # flows.
    _force_heading_autofit(shape, text)
    shape.top = top
    # Only pin height back to the template's fixed single-line value when
    # the text actually still fits on one line. If _force_heading_autofit
    # had to turn wrapping back on (heading too long even after widening
    # the pill to the slide edge), restoring the old fixed height here
    # would re-clip that 2nd line right back outside the shape -- so keep
    # whatever taller height it computed in that case instead.
    if not shape.text_frame.word_wrap:
        shape.height = height
    try:
        shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass


def _first_run_font_size_pt(shape):
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        sz = tf.paragraphs[0].runs[0].font.size
        if sz:
            return sz.pt
    return 72.0  # this template's observed default


def _estimate_theory_capacity(theory_body_shape, slide_height_emu, bottom_margin_in=0.6):
    """Reads the TEMPLATE's real box width/position/font-size and estimates
    how many characters of theory text fit on one slide, so pagination
    adapts automatically if the template geometry ever changes -- never a
    hardcoded 'N characters per slide' guess."""
    font_size_pt = _first_run_font_size_pt(theory_body_shape)
    box_width_in = theory_body_shape.width / 914400
    box_top_in = theory_body_shape.top / 914400
    slide_height_in = slide_height_emu / 914400

    CHAR_WIDTH_FACTOR = 0.52   # empirical avg glyph width, proportional serif font
    LINE_HEIGHT_FACTOR = 1.25
    char_w_in = (font_size_pt / 72.0) * CHAR_WIDTH_FACTOR
    line_h_in = (font_size_pt / 72.0) * LINE_HEIGHT_FACTOR
    chars_per_line = max(10, int(box_width_in / char_w_in))
    avail_height_in = max(line_h_in, slide_height_in - box_top_in - bottom_margin_in)
    max_lines = max(1, int(avail_height_in / line_h_in))
    return chars_per_line * max_lines


def _estimate_text_height_emu(shape, text):
    """Best-effort estimate (EMU) of how tall `text` will render inside
    `shape`'s current width/font-size -- used to catch overlap between
    the question stem and the options below it even when the stem
    shape's own saved height doesn't reflect true overflow (e.g. a
    template box that visually overflows instead of growing)."""
    text = text or ""
    if not text.strip():
        return 0
    # Native Word OMML is carried through the edited-DOCX pipeline as a
    # base64 token.  Never measure the token itself (it can be thousands of
    # characters); measure its visible math representation instead.
    text_for_measure = render_mixed_text(text) if 'OMML_TOKEN_RE' in globals() else text
    font_size_pt = _first_run_font_size_pt(shape)
    box_width_in = max(0.5, shape.width / 914400)
    CHAR_WIDTH_FACTOR = 0.52
    LINE_HEIGHT_FACTOR = 1.25
    char_w_in = (font_size_pt / 72.0) * CHAR_WIDTH_FACTOR
    line_h_in = (font_size_pt / 72.0) * LINE_HEIGHT_FACTOR
    chars_per_line = max(10, int(box_width_in / char_w_in)) if char_w_in > 0 else 40
    total_lines = 0
    for line in text_for_measure.split("\n"):
        ln = len(line)
        total_lines += max(1, -(-ln // chars_per_line))  # ceil division
    total_lines = max(1, total_lines)
    return int(total_lines * line_h_in * 914400)


def _wrap_pack(text, max_chars):
    """Greedy pack up to max_chars, cutting on the last whole-word boundary
    (never mid-word) so a slide never gets a chopped-off word. Also never
    cuts through the MIDDLE of a protected math span ($..$, \\(..\\),
    \\[..\\], \\begin{env}..\\end{env}) -- if the word-boundary cut would
    land inside one, the cut is pulled back to before that span so the
    whole equation moves to the next slide intact instead of being torn
    in half (which used to silently corrupt/garble split equations)."""
    text = text.strip()
    if len(text) <= max_chars:
        return text, ""
    cut = text.rfind(" ", 0, max_chars)
    if cut <= 0:
        cut = max_chars
    for m in MATH_SPAN_PATTERN.finditer(text):
        if m.start() < cut < m.end():
            cut = m.start()
            break
    if cut <= 0:
        cut = max_chars  # no safe word/pre-span boundary found -- last resort
    return text[:cut].strip(), text[cut:].strip()


# --- Segmenting source content into heading / body / question blocks ---
# NOTE (Swapnil, please read): question-boundary detection here is a
# best-effort heuristic (source has no explicit "this is a question"
# flag) -- a block is treated as a NEW question only if it starts with a
# source question-number marker (QNUM_PATTERN, e.g. "12.", "Q3)") AND is
# short (< QUESTION_MAX_CHARS) so long numbered theory paragraphs like
# "1. Introduction to electrostatics..." aren't misread as questions.
# Tune QUESTION_MAX_CHARS below against your real PDFs if you see misfires.
QUESTION_MAX_CHARS = 600


# --- Real equation rendering (matplotlib mathtext -- free, offline, no ----
# --- external LaTeX install needed) for PPTX, where unlike Word there's --
# --- no OMML-in-textbox support via python-pptx. This produces an actual
# --- crisp math image instead of a Unicode approximation. -----------------
_MPL_READY = None


def _mpl_available():
    global _MPL_READY
    if _MPL_READY is None:
        try:
            import matplotlib
            matplotlib.use('Agg')
            _MPL_READY = True
        except Exception:
            _MPL_READY = False
    return _MPL_READY


_MPL_FIG_CACHE = None


def render_latex_to_png_bytes(latex, fontsize=44):
    """Renders LaTeX to a tightly-cropped transparent PNG via matplotlib's
    mathtext engine. Returns (png_bytes, width_px, height_px), or None if
    mathtext can't parse this particular LaTeX (e.g. matrix environments) --
    callers must fall back to latex_to_display_text() in that case so
    nothing is ever silently dropped."""
    global _MPL_FIG_CACHE
    if not latex or not latex.strip() or not _mpl_available():
        return None
    import matplotlib.pyplot as plt
    s = latex.strip()
    s = re.sub(r'^\${1,2}|\${1,2}$', '', s).strip()
    if not s:
        return None
    s = '$' + s + '$'
    try:
        if _MPL_FIG_CACHE is None:
            _MPL_FIG_CACHE = plt.figure(figsize=(0.1, 0.1))
        fig = _MPL_FIG_CACHE
        fig.clf()
        fig.patch.set_alpha(0)
        t = fig.text(0, 0, s, fontsize=fontsize, color='black')
        fig.canvas.draw()
        bbox = t.get_window_extent()
        if bbox.width <= 0 or bbox.height <= 0:
            return None
        fig.set_size_inches(bbox.width / fig.dpi, bbox.height / fig.dpi)
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=300, transparent=True, bbox_inches='tight', pad_inches=0.03)
        buf.seek(0)
        png_bytes = buf.getvalue()
        try:
            from PIL import Image
            img = Image.open(io.BytesIO(png_bytes))
            w, h = img.size
        except Exception:
            w, h = max(1, int(bbox.width * 3)), max(1, int(bbox.height * 3))
        return png_bytes, w, h
    except Exception:
        try:
            if _MPL_FIG_CACHE is not None:
                _MPL_FIG_CACHE.clf()
        except Exception:
            pass
        return None


def _fit_dims(img_w_px, img_h_px, max_w_emu, max_h_emu):
    """Contain-fit: scales (img_w_px, img_h_px) to the largest size that
    fits within (max_w_emu, max_h_emu), preserving aspect ratio."""
    if img_w_px <= 0 or img_h_px <= 0:
        return max_w_emu, max_h_emu
    aspect = img_w_px / img_h_px
    w = max_w_emu
    h = int(w / aspect)
    if h > max_h_emu:
        h = max_h_emu
        w = int(h * aspect)
    return max(1, w), max(1, h)


def _add_picture_fit(slide, image_source, left_emu, top_emu, max_w_emu, max_h_emu, img_w_px=None, img_h_px=None):
    """Adds a picture (image_source: BytesIO or file path) contain-fit and
    CENTERED within the given box, so it never overflows into neighboring
    shapes. Reads pixel dims itself if not supplied."""
    if img_w_px is None or img_h_px is None:
        try:
            from PIL import Image
            if hasattr(image_source, 'read'):
                pos = image_source.tell()
                img = Image.open(image_source)
                img_w_px, img_h_px = img.size
                image_source.seek(pos)
            else:
                img = Image.open(image_source)
                img_w_px, img_h_px = img.size
        except Exception:
            img_w_px, img_h_px = max_w_emu, max_h_emu  # fallback: assume it already fits
    w, h = _fit_dims(img_w_px, img_h_px, max_w_emu, max_h_emu)
    left = left_emu + (max_w_emu - w) // 2
    top = top_emu + (max_h_emu - h) // 2
    if hasattr(image_source, 'seek'):
        image_source.seek(0)
    return slide.shapes.add_picture(image_source, left, top, width=w, height=h)


# NOTE: _replace_shape_with_equation_image() and
# _replace_shape_with_native_equation() below are kept for reference/
# possible reuse elsewhere, but are NO LONGER called by the main PPTX
# build pipeline -- every shape (question stems, MCQ options, theory
# body text, standalone equation slides) now goes through
# _set_shape_text_with_math() instead, which never falls back to a
# static (non-editable) image; a failed OMML conversion degrades to
# editable Unicode-approximated text instead.
def _replace_shape_with_equation_image(slide, shape, latex, fontsize=44):
    """Tries to render `latex` as a real math image and swap it in at the
    exact position/size of `shape` (removing the original textbox).
    Falls back to plain Unicode-approximated text in the SAME shape if
    mathtext can't parse this equation -- so nothing is ever dropped, it
    just degrades from 'real image' to 'best-effort text'. Returns True if
    a real image was used, False if it fell back to text."""
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    result = render_latex_to_png_bytes(latex, fontsize=fontsize)
    if result is None:
        _set_shape_text(shape, latex_to_display_text(latex))
        return False
    png_bytes, w_px, h_px = result
    _remove_shape(shape)
    _add_picture_fit(slide, io.BytesIO(png_bytes), left, top, width, height, w_px, h_px)
    return True


def _add_table_shape(slide, rows, left_emu, top_emu, max_w_emu, max_h_emu):
    """Real editable PowerPoint table with native OMML inside cells. Image
    tokens are overlaid within their corresponding cell. Cell content
    (plain text AND any embedded equations) is stamped at a FIXED 60pt
    font size, per explicit requirement -- table content should always
    render at 60pt regardless of the source PDF's own table font size."""
    if not rows: return None
    n_rows=len(rows); n_cols=max(len(r) for r in rows)
    h=min(max_h_emu,max(int(0.9*914400)*n_rows,int(0.7*914400))); w=max_w_emu
    gshape=slide.shapes.add_table(n_rows,n_cols,left_emu,top_emu,w,h); table=gshape.table
    for ci in range(n_cols): table.columns[ci].width=int(w/n_cols)
    for ri in range(n_rows): table.rows[ri].height=int(h/n_rows)
    table_rPr = etree.Element(_aqn('rPr'))
    table_rPr.set('sz', '6000')  # 6000 = 60pt (OOXML sz is in hundredths of a point)
    overlays=[]
    for ri,row in enumerate(rows):
        for ci in range(n_cols):
            raw=row[ci] if ci<len(row) else ''
            imgs=[]
            clean=IMAGE_TOKEN_RE.sub(lambda m:(imgs.append(m.group(1)) or ''),raw)
            cell=table.cell(ri,ci); tf=cell.text_frame; tx=tf._txBody
            for old in list(tx.p_lst): tx.remove(old)
            for line in (clean.split('\n') or ['']): tx.append(_build_mixed_paragraph_xml(line,None,table_rPr,None))
            overlays += [(ri,ci,x) for x in imgs]
    for ri,ci,relpath in overlays:
        candidates=[relpath,os.path.join(os.getcwd(),relpath)]
        img=next((x for x in candidates if os.path.exists(x)),None)
        if not img: continue
        cleft=left_emu+sum(table.columns[j].width for j in range(ci)); ctop=top_emu+sum(table.rows[j].height for j in range(ri))
        try: _add_picture_fit(slide,img,cleft+45720,ctop+45720,int(table.columns[ci].width*.9),int(table.rows[ri].height*.9))
        except Exception: pass
    return gshape


MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
A14_NS = "http://schemas.microsoft.com/office/drawing/2010/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _aqn(local):
    return '{%s}%s' % (A_NS, local)


def _apply_running_text_format_to_omml(omath_el, template_rPr):
    """Native OMML math runs (m:r) carry NO a:rPr of their own by default
    -- mathml2omml only emits m:rPr (m:sty for italic/plain), never a
    size or color. Without an explicit a:rPr, PowerPoint renders the
    equation with its own default math-object appearance (black text,
    a fixed ~40pt size) instead of inheriting the surrounding run's
    formatting -- this is what made equations look black/oversized
    regardless of the template. Fix: stamp the SAME font size + color as
    the template's running text onto every m:r in the equation, so it
    always matches the paragraph it sits in. Only size + color are
    copied (not bold/italic) so math's own italic-variable styling
    (m:sty) is left alone."""
    if omath_el is None or template_rPr is None:
        return omath_el
    sz = template_rPr.get('sz')
    solid_fill = template_rPr.find(pptx_qn('a:solidFill'))
    if sz is None and solid_fill is None:
        return omath_el
    for r_el in omath_el.iter(f'{{{MATH_NS}}}r'):
        # Drop any pre-existing a:rPr on this run before re-adding ours.
        for old in r_el.findall(_aqn('rPr')):
            r_el.remove(old)
        new_rpr = etree.Element(_aqn('rPr'))
        if sz is not None:
            new_rpr.set('sz', sz)
        if solid_fill is not None:
            new_rpr.append(copy.deepcopy(solid_fill))
        t_el = r_el.find(f'{{{MATH_NS}}}t')
        if t_el is not None:
            t_el.addprevious(new_rpr)
        else:
            r_el.append(new_rpr)
    return omath_el


def _build_equation_paragraph_xml(omath_el, fallback_text, template_pPr=None, template_rPr=None):
    """Builds a real <a:p> DrawingML paragraph containing a NATIVE,
    PowerPoint-editable OMML equation -- the exact same m:oMath format
    Word uses, wrapped in the mc:AlternateContent/a14:m mechanism that
    PowerPoint's own 'Insert > Equation' feature generates. PowerPoint
    2010+ renders/edits the real math (the mc:Choice branch); anything
    older automatically shows the plain-text mc:Fallback instead -- this
    is a standard OOXML compatibility feature, not a workaround."""
    p = etree.Element(_aqn('p'), nsmap={'a': A_NS})
    if template_pPr is not None:
        p.append(copy.deepcopy(template_pPr))
    alt = etree.SubElement(p, '{%s}AlternateContent' % MC_NS, nsmap={'mc': MC_NS})
    choice = etree.SubElement(alt, '{%s}Choice' % MC_NS, nsmap={'a14': A14_NS})
    choice.set('Requires', 'a14')
    a14m = etree.SubElement(choice, '{%s}m' % A14_NS)
    a14m.append(copy.deepcopy(omath_el))
    fallback = etree.SubElement(alt, '{%s}Fallback' % MC_NS)
    r = etree.SubElement(fallback, _aqn('r'))
    if template_rPr is not None:
        r.append(copy.deepcopy(template_rPr))
    t = etree.SubElement(r, _aqn('t'))
    t.text = fallback_text
    return p


def _replace_shape_with_native_equation(slide, shape, latex, img_fontsize=44):
    """PRIMARY equation path for PPTX: inserts a real, PowerPoint-editable
    OMML equation (reusing the SAME converter that already works for
    Word). Degrades gracefully if that fails for this specific equation:
    1) real OMML (editable, native)  2) rendered image (correct-looking,
    not editable)  3) plain Unicode text -- so something always displays
    correctly. Returns 'omml' / 'image' / 'text' indicating which path
    was used."""
    tf = shape.text_frame
    first_p = tf.paragraphs[0]
    template_pPr = first_p._p.find(pptx_qn('a:pPr'))
    template_rPr = first_p.runs[0]._r.find(pptx_qn('a:rPr')) if first_p.runs else None

    omath_el = None
    try:
        omath_el = convert_latex_to_omml_native(latex)
    except Exception:
        omath_el = None

    if omath_el is not None:
        try:
            omath_el = _apply_running_text_format_to_omml(omath_el, template_rPr)
            fallback_text = latex_to_display_text(latex)
            new_p = _build_equation_paragraph_xml(omath_el, fallback_text, template_pPr, template_rPr)
            txBody = tf._txBody
            for old_p in list(txBody.p_lst):
                txBody.remove(old_p)
            txBody.append(new_p)
            return "omml"
        except Exception:
            pass  # fall through to image/text fallback below

    result = render_latex_to_png_bytes(latex, fontsize=img_fontsize)
    if result is not None:
        png_bytes, w_px, h_px = result
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        _remove_shape(shape)
        _add_picture_fit(slide, io.BytesIO(png_bytes), left, top, width, height, w_px, h_px)
        return "image"

    _set_shape_text(shape, latex_to_display_text(latex))
    return "text"


# ======================================================================
# MIXED TEXT+MATH SHAPE FILLER (PPTX) -- turns EVERY equation inside a
# shape's text into a real, PowerPoint-editable OMML equation, inline
# with the surrounding plain text, instead of ever flattening a line to
# Unicode-approximated text or dropping to a static (non-editable)
# image. Mirrors what _add_text_paragraph() already does for Word;
# this is the PPTX equivalent, built at the run level so text and math
# can sit on the SAME line/paragraph.
#   - a math span that converts cleanly -> real inline OMML (editable)
#   - a math span that fails to convert -> Unicode-approx text for just
#     that one span (still editable text, never an image)
#   - plain text -> a normal run
# Used for: question stems, MCQ options, and theory paragraph bodies --
# i.e. every place equations were previously being "flattened" or
# turned into pictures.
# ======================================================================

def _build_inline_equation_run_xml(omath_el, fallback_text, template_rPr=None):
    """Run-level <mc:AlternateContent> for ONE equation embedded mid-
    paragraph, alongside normal <a:r> text runs (same mc:Choice/a14:m
    mechanism PowerPoint's own Insert>Equation uses, just without the
    <a:p> wrapper since this sits as a sibling inside an existing one)."""
    alt = etree.Element('{%s}AlternateContent' % MC_NS, nsmap={'mc': MC_NS})
    choice = etree.SubElement(alt, '{%s}Choice' % MC_NS, nsmap={'a14': A14_NS})
    choice.set('Requires', 'a14')
    a14m = etree.SubElement(choice, '{%s}m' % A14_NS)
    a14m.append(copy.deepcopy(omath_el))
    fallback = etree.SubElement(alt, '{%s}Fallback' % MC_NS)
    r = etree.SubElement(fallback, _aqn('r'))
    if template_rPr is not None:
        r.append(copy.deepcopy(template_rPr))
    t = etree.SubElement(r, _aqn('t'))
    t.text = fallback_text
    return alt


def _build_mixed_paragraph_xml(line_text, template_pPr=None, template_rPr=None, math_stats=None, force_align=None):
    """Builds one <a:p> for a single line of source text that may contain
    inline math spans. Plain text -> normal <a:r> runs. Each math span ->
    a real editable OMML equation (via convert_latex_to_omml_native, the
    SAME converter already used for Word/standalone-equation slides).
    If OMML conversion fails for one particular span, ONLY that span
    degrades to Unicode-approximated text -- everything else on the line,
    including other equations, stays fully editable. Never produces an
    image. math_stats, if given, is a dict this increments in place with
    'omml' / 'text_fallback' counts for reporting.

    force_align, if given ('l'/'ctr'/'r'), is stamped as the paragraph's
    a:pPr/@algn attribute DIRECTLY on the XML we build here -- not applied
    afterwards via python-pptx. This matters specifically for a paragraph
    whose ONLY content is an equation (an MCQ option that's pure math):
    PowerPoint's own math-object rendering can visually re-center such a
    paragraph even when a post-hoc `paragraph.alignment = LEFT` has been
    set on the same pPr element, because the alignment must already be
    present on the pPr the very first time PowerPoint's math layout pass
    reads it. Baking it in at construction time removes that ambiguity."""
    p = etree.Element(_aqn('p'), nsmap={'a': A_NS})
    pPr_el = None
    if template_pPr is not None:
        pPr_el = copy.deepcopy(template_pPr)
        p.append(pPr_el)
    if force_align is not None:
        if pPr_el is None:
            pPr_el = etree.Element(_aqn('pPr'))
            p.insert(0, pPr_el)
        pPr_el.set('algn', force_align)

    def _add_text_run(content, bold=False, italic=False):
        if content == "":
            return
        r = etree.SubElement(p, _aqn('r'))
        rPr = copy.deepcopy(template_rPr) if template_rPr is not None else None
        if bold or italic:
            if rPr is None:
                rPr = etree.Element(_aqn('rPr'))
            if bold:
                rPr.set('b', '1')
            if italic:
                rPr.set('i', '1')
        if rPr is not None:
            r.append(rPr)
        t = etree.SubElement(r, _aqn('t'))
        t.text = content

    wrote_any = False
    pieces = list(_iter_text_math_or_omml(line_text))
    has_real_text = any(
        piece[0] == "text" and (piece[1] or "").strip()
        for piece in pieces
    )
    # ADVANCED FIX for pure-equation paragraphs (e.g. an MCQ option that is
    # ENTIRELY one equation, no leading/trailing text): even with algn="l"
    # baked into pPr (force_align above), PowerPoint's math layout engine
    # can still visually re-center a paragraph whose ONLY content is a
    # math object -- this is a PowerPoint rendering quirk tied to there
    # being literally nothing else on the line, not to the pPr alignment
    # value itself. The reliable fix (confirmed by hand in PowerPoint): put
    # one invisible leading text run BEFORE the equation, so the paragraph
    # is no longer "pure math" from PowerPoint's point of view. A single
    # space is used (rather than a zero-width character) because it has
    # universal font support; at normal MCQ-option font sizes the visual
    # shift is a fraction of a point and not perceptible, while the
    # left-alignment now holds reliably.
    if force_align is not None and pieces and not has_real_text:
        _add_text_run(" ")
        wrote_any = True

    for piece in pieces:
        kind = piece[0]
        if kind == "text":
            _, content, bold, italic = piece
            if content:
                _add_text_run(content, bold, italic)
                wrote_any = True
            continue
        if kind == "omml":
            omath_el = piece[1]
            if omath_el is None:
                continue
            try:
                omath_el = _apply_running_text_format_to_omml(omath_el, template_rPr)
                p.append(_build_inline_equation_run_xml(omath_el, _omml_visible_text(omath_el), template_rPr))
                wrote_any = True
                if math_stats is not None:
                    math_stats["omml"] = math_stats.get("omml", 0) + 1
            except Exception:
                _add_text_run(_omml_visible_text(omath_el))
                wrote_any = True
                if math_stats is not None:
                    math_stats["text_fallback"] = math_stats.get("text_fallback", 0) + 1
            continue
        # kind == "math"
        content = piece[1]
        if not (content or "").strip():
            continue
        omath_el = None
        try:
            omath_el = convert_latex_to_omml_native(content)
        except Exception:
            omath_el = None
        if omath_el is not None:
            try:
                omath_el = _apply_running_text_format_to_omml(omath_el, template_rPr)
                p.append(_build_inline_equation_run_xml(omath_el, latex_to_display_text(content), template_rPr))
                wrote_any = True
                if math_stats is not None:
                    math_stats["omml"] = math_stats.get("omml", 0) + 1
                continue
            except Exception:
                pass
        _add_text_run(latex_to_display_text(content))
        wrote_any = True
        if math_stats is not None:
            math_stats["text_fallback"] = math_stats.get("text_fallback", 0) + 1

    if not wrote_any:
        _add_text_run("")
    return p


def _set_shape_text_with_math(shape, raw_text, math_stats=None, force_align=None):
    """Drop-in replacement for _set_shape_text() that ALSO converts every
    inline LaTeX math span in raw_text into a real, PowerPoint-editable
    OMML equation, instead of flattening the whole line to plain text.
    Clones the template shape's original paragraph/run formatting exactly
    like _set_shape_text() did, so styling (font, size, color, alignment)
    is unaffected.

    force_align ('l'/'ctr'/'r'), if given, forces every paragraph's algn
    explicitly -- see _build_mixed_paragraph_xml for why this is baked in
    at construction time rather than patched afterwards."""
    tf = shape.text_frame
    first_p = tf.paragraphs[0]
    pPr_src = first_p._p.find(pptx_qn('a:pPr'))
    rPr_src = None
    if first_p.runs:
        rPr_src = first_p.runs[0]._r.find(pptx_qn('a:rPr'))

    txBody = tf._txBody
    for old_p in list(txBody.p_lst):
        txBody.remove(old_p)

    lines = (raw_text or "").split("\n") or [""]
    for line in lines:
        txBody.append(_build_mixed_paragraph_xml(line, pPr_src, rPr_src, math_stats, force_align=force_align))


def _opt_marker_index(display_text):
    """Canonical 1..4 index for an option's leading marker (A/1->1 ...
    D/4->4), or None if the text doesn't start with one."""
    m = OPT_PATTERN.match((display_text or "").strip())
    if not m:
        return None
    marker = re.sub(r'[^A-Za-z0-9]', '', m.group(1)).upper()
    return {'A': 1, 'B': 2, 'C': 3, 'D': 4, '1': 1, '2': 2, '3': 3, '4': 4}.get(marker)


def segment_content_for_pptx(content_list_path):
    with open(content_list_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    pages = {}
    for item in data:
        if item.get('type') in PAGE_FURNITURE_TYPES:
            continue
        pages.setdefault(item.get('page_idx', 0), []).append(item)

    segments = []
    q_buffer = []  # list of dicts: {"item_type","display","raw"}
    in_question = [False]

    def flush_question():
        if not q_buffer:
            in_question[0] = False
            return
        text_entries = [e for e in q_buffer if e["item_type"] in ("text", "equation")]
        extra_blocks = [e for e in q_buffer if e["item_type"] in ("image", "table")]

        opt_positions = [i for i, e in enumerate(text_entries)
                          if OPT_PATTERN.match(e["display"].strip())]
        if len(opt_positions) >= 2:
            stem = " ".join(e["display"] for e in text_entries[:opt_positions[0]]).strip()
            stem_raw = " ".join(e["raw"] for e in text_entries[:opt_positions[0]]).strip()
            options = []
            for i in opt_positions[:4]:
                entry = text_entries[i]
                marker, rest = split_leading_marker(entry["display"])
                option_display = rest if rest else entry["display"]
                # Every option carries its RAW source text (with math
                # delimiters intact) so the renderer can turn any inline
                # equation in it into a real, editable OMML equation --
                # no more all-or-nothing "is this whole option basically
                # one equation" heuristic, and never a static image.
                _, raw_rest = split_leading_marker(entry["raw"])
                option_raw = raw_rest if raw_rest else entry["raw"]
                options.append({"display": option_display, "raw": option_raw})
            extra_note = ""
            if len(opt_positions) > 4:
                extra_note = (" [%d extra option-like line(s) beyond D were found and "
                               "NOT shown -- manual review]" % (len(opt_positions) - 4))
            segments.append({"kind": "question", "stem": (stem + extra_note).strip(),
                              "stem_raw": stem_raw, "options": options, "qtype": "mcq",
                              "extra_blocks": extra_blocks})
        else:
            stem = " ".join(e["display"] for e in text_entries).strip()
            stem_raw = " ".join(e["raw"] for e in text_entries).strip()
            segments.append({"kind": "question", "stem": stem, "stem_raw": stem_raw,
                              "options": [], "qtype": "other", "extra_blocks": extra_blocks})
        q_buffer.clear()
        in_question[0] = False

    def _current_opt_count():
        return sum(1 for e in q_buffer
                   if e["item_type"] in ("text", "equation") and OPT_PATTERN.match(e["display"].strip()))

    def _maybe_split_on_option_restart():
        """Detects a SECOND question's options starting (marker cycles
        back to A/1) inside the SAME buffer -- this is what used to let
        two MCQs silently land on one slide when the source gave no clean
        separator between "...D) ..." and the next question's own
        "A) ...". As soon as an option marker resets after 4 real options
        have already been seen, everything gathered so far is flushed as
        its own complete question and a fresh buffer picks up right at
        the reset point (still counted as "in a question", so the new
        question's own stem/options keep accumulating normally)."""
        idxs = [i for i, e in enumerate(q_buffer)
                if e["item_type"] in ("text", "equation") and OPT_PATTERN.match(e["display"].strip())]
        if len(idxs) < 5:
            return
        seen = 0
        last = 0
        last_option_buf_idx = None
        for idx in idxs:
            mi = _opt_marker_index(q_buffer[idx]["display"])
            if mi is None:
                continue
            if mi <= last and seen >= 4:
                # Split right AFTER the previous question's 4th real
                # option -- not at the reset option itself -- so any
                # stray non-option line sitting between (almost always
                # the next question's own stem, since a new question
                # always needs one before its options) is carried into
                # the new question's buffer instead of being silently
                # dropped during the completed question's flush.
                split_at = last_option_buf_idx + 1
                tail = q_buffer[split_at:]
                del q_buffer[split_at:]
                flush_question()
                q_buffer.extend(tail)
                in_question[0] = True
                return
            last = mi
            seen += 1
            last_option_buf_idx = idx

    for page_idx in sorted(pages.keys()):
        for item in _order_page_items_by_column(pages[page_idx]):
            itype = item.get('type', 'text')
            text_level = item.get('text_level', 0) or 0

            # FIX: previously ANY text_level >= 1 (main heading AND every
            # sub-heading under it) was treated as an equal "heading"
            # segment, and each one simply OVERWROTE current_heading --
            # so when a main heading was immediately followed by its own
            # sub-heading (e.g. "Clothing & Ornaments" -> "Clothing"),
            # only the LAST one before body text survived into the
            # slide's single Heading placeholder. The main heading was
            # silently discarded and the sub-heading got promoted into
            # its place. Only a genuine top-level heading (itype=='title'
            # or text_level == 1) may now replace the slide heading;
            # anything deeper (text_level >= 2) is kept as bold body text
            # instead, so it stays visible under the CORRECT heading
            # rather than replacing it.
            if itype == 'title' or text_level == 1:
                flush_question()
                segments.append({"kind": "heading", "text": (item.get('text') or '').strip()})
                continue

            if text_level >= 2:
                flush_question()
                sub_text = (item.get('text') or '').strip()
                if sub_text:
                    segments.append({"kind": "text", "text": render_mixed_text("**" + sub_text + "**"),
                                      "raw": "**" + sub_text + "**"})
                continue

            if itype == 'equation':
                raw_latex = _strip_latex_delimiters((item.get('text') or '').strip())
                if not raw_latex:
                    continue
                conv = latex_to_display_text(raw_latex)
                if in_question[0]:
                    # Wrap in $..$ so this entry's "raw" is consistent
                    # with every other entry's raw (plain text that may
                    # CONTAIN a delimited math span) -- one uniform
                    # format for the mixed text+math renderer downstream.
                    q_buffer.append({"item_type": "equation", "display": conv,
                                      "raw": "$" + raw_latex + "$"})
                    _maybe_split_on_option_restart()
                    if _current_opt_count() >= 8:
                        flush_question()
                else:
                    segments.append({"kind": "equation", "latex": raw_latex})
                continue

            if itype == 'table':
                table_html = item.get('table_body', '')
                rows = parse_html_table_to_rows(table_html) if table_html else []
                if in_question[0]:
                    q_buffer.append({"item_type": "table", "rows": rows})
                else:
                    if rows:
                        segments.append({"kind": "table", "rows": rows})
                    else:
                        raw_fallback = item.get('text') or ''
                        if raw_fallback.strip():
                            segments.append({"kind": "text", "text": render_mixed_text(raw_fallback),
                                              "raw": raw_fallback})
                continue

            if itype in ('image', 'chart'):
                img_path = item.get('img_path')
                if in_question[0]:
                    q_buffer.append({"item_type": "image", "img_path": img_path, "page_idx": page_idx})
                else:
                    segments.append({"kind": "image", "img_path": img_path, "page_idx": page_idx})
                continue

            raw_text = (item.get('text', '') or '').strip()
            if not raw_text:
                continue

            # Un-glue any options bunched onto one line (e.g. "A) x B) y
            # C) z D) w" all in a single string) onto their own vertical
            # lines -- the SAME normalization the DOCX pipeline already
            # relies on (process_mcq_formatting), just not previously
            # applied here.
            normalized = process_mcq_formatting(raw_text)
            norm_lines = normalized.split("\n")
            has_qnum_inside = any(QNUM_PATTERN.match(ln.strip()) for ln in norm_lines)

            if not in_question[0] and not has_qnum_inside:
                # Pure flowing theory content, no question markers at all
                # -- keep the item as ONE unit so its internal line breaks
                # (paragraph structure) are preserved exactly as extracted.
                display_text = render_mixed_text(normalized)
                segments.append({"kind": "text", "text": display_text, "raw": normalized})
                continue

            # Otherwise this item is (or contains) question content --
            # walk it line by line so EVERY question-number marker and
            # EVERY option marker is evaluated independently. This is
            # what lets a second (or third) question glued into the same
            # extracted block -- with no page/paragraph break MinerU
            # could use to separate them -- still get its own slide
            # instead of silently merging into the previous question.
            for line in norm_lines:
                line = line.strip()
                if not line:
                    continue
                line_display = render_mixed_text(line)
                starts_question = (bool(QNUM_PATTERN.match(line_display))
                                    and len(line_display) < QUESTION_MAX_CHARS)

                if starts_question:
                    flush_question()
                    in_question[0] = True
                    # Strip the source question-number marker (e.g. "1.",
                    # "Q3)") BEFORE buffering -- otherwise OPT_PATTERN's
                    # digit-marker branch (which matches "1)"/"2)" style
                    # options) collides with the question number itself
                    # and misreads it as option 1, wiping the real stem.
                    qnum_m = QNUM_PATTERN.match(line_display)
                    disp = line_display[qnum_m.end():].strip() if qnum_m else line_display
                    raw_stripped = QNUM_PATTERN.sub('', line, count=1).strip()
                    q_buffer.append({"item_type": "text", "display": disp, "raw": raw_stripped})
                elif in_question[0]:
                    q_buffer.append({"item_type": "text", "display": line_display, "raw": line})
                    _maybe_split_on_option_restart()
                    if _current_opt_count() >= 8:
                        flush_question()
                else:
                    segments.append({"kind": "text", "text": line_display, "raw": line})

    flush_question()
    return segments



_WORD_NS = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'


def _word_run_is_bold(r_el):
    rpr = r_el.find(f'{{{_WORD_NS}}}rPr')
    if rpr is None:
        return False
    b = rpr.find(f'{{{_WORD_NS}}}b')
    if b is None:
        return False
    val = b.get(f'{{{_WORD_NS}}}val')
    return val not in ('0', 'false', 'off')


def _word_run_is_italic(r_el):
    rpr = r_el.find(f'{{{_WORD_NS}}}rPr')
    if rpr is None:
        return False
    i = rpr.find(f'{{{_WORD_NS}}}i')
    if i is None:
        return False
    val = i.get(f'{{{_WORD_NS}}}val')
    return val not in ('0', 'false', 'off')


# FIX 4 (Swapnil, theory images missing): images that were manually
# pasted into the edited Word doc (as opposed to ones this app itself
# inserted via python-docx) sometimes land as legacy VML drawings
# (<w:pict><v:shape><v:imagedata r:id="..."/>) instead of the modern
# DrawingML <w:drawing><...><a:blip r:embed="..."/> this code previously
# looked for exclusively -- Word still writes VML for pictures pasted
# via some "Paste Special" options, for content copied out of older
# .doc files, and for certain compatibility-mode edits. Those images
# were silently dropped because only 'a:blip' was ever checked. Now
# BOTH the modern blip AND legacy v:imagedata markers are checked (VML
# uses either the modern r:id attribute or the legacy o:relid one), and
# a run is never double-counted -- if a modern blip was already found,
# the VML fallback (which some AlternateContent blocks carry alongside
# the modern version for backward compatibility) is skipped.
_R_NS = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'


def _word_run_image_paths(run_el, rel_to_path):
    """Every image (modern DrawingML blip OR legacy VML imagedata) found
    inside this run element, resolved to an extracted file path via
    rel_to_path. Never returns the same image twice for one run."""
    found = []
    seen_rids = set()
    for el in run_el.iter():
        if el.tag.endswith('}blip'):
            rid = el.get(_R_NS + 'embed') or el.get(_R_NS + 'link')
            if rid and rid not in seen_rids and rid in rel_to_path:
                seen_rids.add(rid)
                found.append(rel_to_path[rid])
    if found:
        return found
    for el in run_el.iter():
        if el.tag.endswith('}imagedata'):
            rid = el.get(_R_NS + 'id') or el.get(_R_NS + 'link') or el.get('o:relid')
            if rid and rid not in seen_rids and rid in rel_to_path:
                seen_rids.add(rid)
                found.append(rel_to_path[rid])
    return found


def _word_paragraph_runs_lossless(paragraph, rel_to_path):
    """Walk a Word paragraph in XML order and preserve native OMML and images.
    Also preserves each run's real bold/italic formatting (w:b / w:i) by
    wrapping its text in **markdown** / *markdown* markers -- the SAME
    convention the PPTX text-run builder (_split_markdown_style) already
    understands for MinerU-sourced content -- so bold/italic applied by
    hand in the edited Word doc survives all the way into the PPTX
    instead of being silently discarded (previously only the plain run
    text was kept, dropping every w:b/w:i flag)."""
    segs=[]
    for child in paragraph._p.iterchildren():
        tag=child.tag
        if tag.endswith('}oMath') or tag.endswith('}oMathPara'):
            token=_encode_omml_token(child)
            if token: segs.append(('omml',token))
            continue
        if tag.endswith('}r'):
            texts=[t.text or '' for t in child.iter() if t.tag=='{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t']
            if texts:
                run_text = ''.join(texts)
                if run_text.strip():
                    is_b = _word_run_is_bold(child)
                    is_i = _word_run_is_italic(child)
                    if is_b and is_i:
                        run_text = '***' + run_text + '***'
                    elif is_b:
                        run_text = '**' + run_text + '**'
                    elif is_i:
                        run_text = '*' + run_text + '*'
                segs.append(('text',run_text))
            for img_val in _word_run_image_paths(child, rel_to_path):
                segs.append(('image', img_val))
            continue
        for img_val in _word_run_image_paths(child, rel_to_path):
            segs.append(('image', img_val))
    return segs


def _word_paragraph_has_bullet(paragraph):
    """True if the Word paragraph itself carries list/bullet numbering
    (w:pPr/w:numPr) -- i.e. a real Word bulleted-list paragraph, as
    opposed to a literal '-'/'*' character typed into plain text. Word
    bullets don't exist as text at all, so without this check they
    vanished completely once re-exported (no bullet was ever "lost" as
    a character -- it simply was never extracted)."""
    pPr = paragraph._p.find(f'{{{_WORD_NS}}}pPr')
    if pPr is None:
        return False
    return pPr.find(f'{{{_WORD_NS}}}numPr') is not None

def _word_paragraph_text_with_math(paragraph):
    parts=[]
    for kind,val in _word_paragraph_runs_lossless(paragraph,{}):
        if kind=='text': parts.append(val)
        elif kind=='omml': parts.append(_omml_visible_text(_decode_omml_token(val)))
    return ''.join(parts).replace('\t','    ').strip()

def _word_paragraph_is_heading(paragraph, text):
    style_name=''
    try: style_name=(paragraph.style.name or '').lower()
    except Exception: pass
    if style_name.startswith('heading') or style_name in ('title','subtitle'): return True
    runs=[r for r in paragraph.runs if (r.text or '').strip()]
    if not runs or not text or len(text)>220: return False
    sizes=[r.font.size.pt for r in runs if r.font.size is not None]
    return all(bool(r.bold) for r in runs) and bool(sizes) and max(sizes)>=13

def _word_extract_images(doc, workspace):
    """Extract every embedded image relationship from the whole DOCX package.

    FIX (Swapnil): the old dedup check compared a (partname, content_type)
    tuple against rel_to_path's KEYS -- but rel_to_path is keyed by rel_id,
    not by (partname, content_type), so that check could never actually
    match anything. Net effect: doc.part is walked once explicitly AND
    again via the package.parts loop below (doc.part's own partname also
    starts with '/word/'), so every image got extracted TWICE under two
    different counter numbers, and rel_to_path silently ended up pointing
    at whichever extraction ran last -- wasteful but usually not
    outright missing. Now dedup is tracked correctly by the target
    part's identity, so each embedded image is written to disk exactly
    once and EVERY relationship id that points at it (there can be more
    than one rId for the same image part) still resolves to that one
    file."""
    img_dir=os.path.join(workspace,'edited_word_media'); os.makedirs(img_dir,exist_ok=True)
    rel_to_path={}; counter=0
    target_to_path={}  # (partname, ctype) -> already-written file path
    parts=[doc.part]
    try:
        parts += [p for p in doc.part.package.parts if getattr(p,'partname',None) and str(p.partname).startswith('/word/')]
    except Exception: pass
    for part in parts:
        for rel_id,rel in getattr(part,'rels',{}).items():
            target=getattr(rel,'target_part',None); ctype=getattr(target,'content_type','') if target is not None else ''
            if not ctype or not ctype.startswith('image/'): continue
            key=(getattr(target,'partname',None),ctype)
            if key in target_to_path:
                rel_to_path[rel_id]=target_to_path[key]
                continue
            try:
                ext=ctype.split('/',1)[1].lower().replace('jpeg','jpg').replace('svg+xml','svg')
                counter+=1; out=os.path.join(img_dir,f'image_{counter}.{ext}')
                with open(out,'wb') as f: f.write(target.blob)
                rel_to_path[rel_id]=out
                target_to_path[key]=out
            except Exception: pass
    return rel_to_path

def _word_cell_raw(cell, rel_to_path, workspace):
    chunks=[]
    for para in cell.paragraphs:
        line=''
        for kind,val in _word_paragraph_runs_lossless(para,rel_to_path):
            if kind in ('text','omml'): line+=val
            elif kind=='image': line+='@@IMG:'+os.path.abspath(val).replace('\\','/')+'@@'
        if line or not chunks: chunks.append(line)
    return '\n'.join(chunks).strip()

def _word_table_to_rows(table, rel_to_path, workspace):
    return [[_word_cell_raw(cell,rel_to_path,workspace) for cell in row.cells] for row in table.rows]

def _word_block_sequence(docx_path, workspace):
    """Lossless edited-DOCX walker. Equations remain native OMML tokens, images
    remain relationship-backed files, and tables remain structured blocks."""
    doc=docx.Document(docx_path)
    rel_to_path=_word_extract_images(doc,workspace)
    body=doc.element.body
    para_map={p._p:p for p in doc.paragraphs}; table_map={t._tbl:t for t in doc.tables}
    blocks=[]; page_idx=0
    for child in body.iterchildren():
        tag=child.tag
        if tag.endswith('}sectPr'): continue
        if tag.endswith('}p'):
            para=para_map.get(child)
            if para is None: continue
            segs=_word_paragraph_runs_lossless(para,rel_to_path)
            # A real Word bulleted-list paragraph (w:numPr) has no bullet
            # CHARACTER in its text at all -- prefix one so the source's
            # list structure survives into the PPTX instead of silently
            # disappearing (previously nothing marked this paragraph as
            # a list item once its text was extracted).
            raw = '\u2022 ' if _word_paragraph_has_bullet(para) else ''
            emitted=False
            for kind,val in segs:
                if kind in ('text','omml'): raw+=val
                elif kind=='image':
                    if raw.strip():
                        visible=render_mixed_text(raw)
                        typ='title' if _word_paragraph_is_heading(para,visible) else 'text'
                        blocks.append({'type':typ,'text':raw if typ=='title' else process_mcq_formatting(raw),'page_idx':page_idx}); page_idx+=1; raw=''; emitted=True
                    blocks.append({'type':'image','img_path':os.path.relpath(val,workspace).replace('\\','/'),'page_idx':page_idx}); page_idx+=1; emitted=True
            if raw.strip():
                visible=render_mixed_text(raw)
                typ='title' if _word_paragraph_is_heading(para,visible) else 'text'
                blocks.append({'type':typ,'text':raw if typ=='title' else process_mcq_formatting(raw),'page_idx':page_idx}); page_idx+=1
            elif not emitted and para.text.strip():
                blocks.append({'type':'text','text':para.text.strip(),'page_idx':page_idx}); page_idx+=1
        elif tag.endswith('}tbl'):
            table=table_map.get(child)
            if table is None: continue
            rows=_word_table_to_rows(table,rel_to_path,workspace)
            if rows:
                import html as _html
                html_rows=[]
                for row in rows:
                    html_rows.append('<tr>'+''.join('<td>'+_html.escape(c).replace('@@OMML:','@@OMML:').replace('@@IMG:','@@IMG:').replace('\n','<br>')+'</td>' for c in row)+'</tr>')
                blocks.append({'type':'table','table_body':'<table>'+''.join(html_rows)+'</table>','page_idx':page_idx}); page_idx+=1
    return blocks,workspace

def build_pptx_from_word(docx_path, template_path, output_pptx_path, workspace=None):
    """Direct edited-WORD -> PPTX path.

    The edited .docx is now the source of truth.  We deliberately do not call
    MinerU, do not rebuild content from the original PDF, and do not require
    content_list.json.  This is what makes Word formatting/content corrections
    actually reach the final PPTX.
    """
    workspace = workspace or os.path.dirname(output_pptx_path)
    os.makedirs(workspace, exist_ok=True)
    blocks, base_dir = _word_block_sequence(docx_path, workspace)
    temp_json = os.path.join(workspace, '_edited_word_content_list.json')
    with open(temp_json, 'w', encoding='utf-8') as f:
        json.dump(blocks, f, ensure_ascii=False, indent=2)
    return build_pptx_from_content_list(temp_json, base_dir, template_path, output_pptx_path)

def build_pptx_from_content_list(content_list_path, base_dir, template_path, output_pptx_path):
    prs = Presentation(template_path)
    if len(prs.slides) < 2:
        raise ValueError(
            "Template me kam se kam 2 slides chahiye (1st = Theory blueprint "
            "with 'TextBox 1', 2nd = Q&A blueprint with 'TextBox 21..25')."
        )
    theory_bp_idx, qna_bp_idx = 0, 1
    theory_blueprint = prs.slides[theory_bp_idx]
    qna_blueprint = prs.slides[qna_bp_idx]

    theory_body_bp_shape = _find_shape(theory_blueprint, THEORY_BODY_SHAPE)
    if theory_body_bp_shape is None:
        raise ValueError(f"Template ke pehle slide me '{THEORY_BODY_SHAPE}' naam ka shape nahi mila.")
    qna_stem_bp_shape = _find_shape(qna_blueprint, Q_STEM_SHAPE)
    if qna_stem_bp_shape is None:
        raise ValueError(f"Template ke doosre slide me '{Q_STEM_SHAPE}' naam ka shape nahi mila.")

    # Legacy character capacity is retained only as a fallback estimate.
    # Actual pagination below is driven by the rendered line-height estimate
    # against the template's REAL expanded TextBox 1 height.
    max_chars_per_theory_slide = _estimate_theory_capacity(theory_body_bp_shape, prs.slide_height)
    theory_content_box = (theory_body_bp_shape.left, theory_body_bp_shape.top,
                           theory_body_bp_shape.width, theory_body_bp_shape.height)

    # Work out the free whitespace gap between the question stem and the
    # first option box, from the TEMPLATE's own geometry -- this is where
    # any diagram/table embedded WITHIN a question gets placed.
    first_opt_bp_shape = _find_shape(qna_blueprint, Q_OPTION_SHAPES[0])
    gap_left = qna_stem_bp_shape.left
    gap_width = qna_stem_bp_shape.width
    gap_top = qna_stem_bp_shape.top + qna_stem_bp_shape.height + int(0.3 * 914400)
    if first_opt_bp_shape is not None:
        gap_bottom = first_opt_bp_shape.top - int(0.3 * 914400)
    else:
        gap_bottom = gap_top + int(3 * 914400)
    gap_height = max(int(0.5 * 914400), gap_bottom - gap_top)

    segments = segment_content_for_pptx(content_list_path)

    stats = {"theory_slides": 0, "mcq_slides": 0, "other_q_slides": 0,
             "options_short": 0, "headings": 0, "equations_as_omml": 0,
             "equations_as_image": 0, "equations_as_text_fallback": 0,
             "images_embedded": 0,
             "tables_embedded": 0, "images_missing": 0}

    def resolve_img(img_path):
        if not img_path:
            return None
        full = os.path.join(base_dir, img_path) if base_dir else img_path
        return full if os.path.exists(full) else None

    def _accumulate_math_stats(math_stats):
        stats["equations_as_omml"] += math_stats.get("omml", 0)
        stats["equations_as_text_fallback"] += math_stats.get("text_fallback", 0)

    def new_theory_text_slide(heading_text, body_text_raw):
        dest = duplicate_slide(prs, theory_bp_idx)
        h_shape = _find_shape(dest, THEORY_HEADING_SHAPE)
        b_shape = _find_shape(dest, THEORY_BODY_SHAPE)
        if h_shape is not None:
            _set_heading_text(h_shape, heading_text or "Notes")
        if b_shape is not None:
            # TextBox 1 in the template is often only one-line tall.  For a
            # theory slide it must become the REAL available body region;
            # otherwise PowerPoint lets text continue below the slide.
            b_shape.height = max(b_shape.height, prs.slide_height - b_shape.top - int(0.55 * 914400))
            math_stats = {}
            _set_shape_text_with_math(b_shape, body_text_raw, math_stats)
            _accumulate_math_stats(math_stats)
        stats["theory_slides"] += 1
        return dest, b_shape

    def new_theory_equation_slide(heading_text, latex):
        dest = duplicate_slide(prs, theory_bp_idx)
        h_shape = _find_shape(dest, THEORY_HEADING_SHAPE)
        b_shape = _find_shape(dest, THEORY_BODY_SHAPE)
        if h_shape is not None:
            _set_heading_text(h_shape, heading_text or "Notes")
        if b_shape is not None:
            # A standalone equation slide is still rendered through the
            # SAME mixed text+math path (never the static-image fallback)
            # so it stays a real, editable OMML equation -- 'editable
            # only', never a picture.
            math_stats = {}
            _set_shape_text_with_math(b_shape, "$" + latex + "$", math_stats)
            _accumulate_math_stats(math_stats)
        stats["theory_slides"] += 1

    def new_theory_image_slide(heading_text, img_path):
        dest = duplicate_slide(prs, theory_bp_idx)
        h_shape = _find_shape(dest, THEORY_HEADING_SHAPE)
        b_shape = _find_shape(dest, THEORY_BODY_SHAPE)
        if h_shape is not None:
            _set_heading_text(h_shape, heading_text or "Notes")
        full_path = resolve_img(img_path)
        if b_shape is not None:
            left, top, width, height = b_shape.left, b_shape.top, b_shape.width, b_shape.height
            # Theory diagrams deserve real vertical room -- extend downward
            # toward the slide bottom rather than being squeezed into
            # TextBox1's normal single-line height.
            height = max(height, prs.slide_height - top - int(1 * 914400))
            _remove_shape(b_shape)
            if full_path:
                _add_picture_fit(dest, full_path, left, top, width, height)
                stats["images_embedded"] += 1
            else:
                # Recreate a small text note so it's still visible that
                # something was here, rather than a silently blank slide.
                tb = dest.shapes.add_textbox(left, top, width, int(0.6 * 914400))
                tb.text_frame.text = "[Image file missing in extracted output -- manual review]"
                stats["images_missing"] += 1
        stats["theory_slides"] += 1

    def new_theory_table_slide(heading_text, rows, leading_text_raw=""):
        """Builds the slide for a table found in the flowing theory text.
        If there's buffered paragraph text that led into this table, it's
        kept on THIS SAME slide (shrunk into the top portion of the body
        box, table placed right below it) instead of being flushed onto
        its own separate slide and pushing the table onto the next one --
        which is what previously made a table look like it "jumped to
        the next page" away from the paragraph introducing it."""
        dest = duplicate_slide(prs, theory_bp_idx)
        h_shape = _find_shape(dest, THEORY_HEADING_SHAPE)
        b_shape = _find_shape(dest, THEORY_BODY_SHAPE)
        if h_shape is not None:
            _set_heading_text(h_shape, heading_text or "Notes")
        if b_shape is not None:
            left, top, width, height = b_shape.left, b_shape.top, b_shape.width, b_shape.height
            full_height = max(height, prs.slide_height - top - int(1 * 914400))
            leading_text_raw = (leading_text_raw or "").strip()
            if leading_text_raw:
                text_h = min(int(full_height * 0.42),
                             max(int(1.1 * 914400), full_height - int(2 * 914400)))
                b_shape.height = text_h
                math_stats = {}
                _set_shape_text_with_math(b_shape, leading_text_raw, math_stats)
                _accumulate_math_stats(math_stats)
                table_top = top + text_h + int(0.15 * 914400)
                table_h = max(int(0.9 * 914400), full_height - text_h - int(0.15 * 914400))
                _add_table_shape(dest, rows, left, table_top, width, table_h)
            else:
                _remove_shape(b_shape)
                _add_table_shape(dest, rows, left, top, width, full_height)
            stats["tables_embedded"] += 1
        stats["theory_slides"] += 1

    def place_extra_blocks(dest, extra_blocks):
        """Places any image/table embedded WITHIN a question into the
        whitespace gap between the stem and the first option box."""
        if not extra_blocks:
            return
        n = len(extra_blocks)
        slot_h = gap_height // n
        for i, block in enumerate(extra_blocks):
            slot_top = gap_top + i * slot_h
            if block["item_type"] == "image":
                full_path = resolve_img(block.get("img_path"))
                if full_path:
                    _add_picture_fit(dest, full_path, gap_left, slot_top, gap_width, slot_h)
                    stats["images_embedded"] += 1
                else:
                    tb = dest.shapes.add_textbox(gap_left, slot_top, gap_width, int(0.6 * 914400))
                    tb.text_frame.text = "[Image file missing in extracted output -- manual review]"
                    stats["images_missing"] += 1
            elif block["item_type"] == "table":
                if block.get("rows"):
                    _add_table_shape(dest, block["rows"], gap_left, slot_top, gap_width, slot_h)
                    stats["tables_embedded"] += 1

    def new_question_slide(seg):
        dest = duplicate_slide(prs, qna_bp_idx)
        stem_shape = _find_shape(dest, Q_STEM_SHAPE)
        stem_raw = seg.get("stem_raw") or seg.get("stem") or "[Question text could not be recovered]"
        if stem_shape is not None:
            math_stats = {}
            _set_shape_text_with_math(stem_shape, stem_raw, math_stats)
            _accumulate_math_stats(math_stats)

        place_extra_blocks(dest, seg.get("extra_blocks") or [])

        if seg["qtype"] == "mcq":
            for i, opt_shape_name in enumerate(Q_OPTION_SHAPES):
                shp = _find_shape(dest, opt_shape_name)
                if shp is None:
                    continue
                if i < len(seg["options"]):
                    opt = seg["options"][i]
                    # Every option -- whether it's plain text, a mix of
                    # text and math, or entirely one equation -- goes
                    # through the SAME editable OMML-first renderer.
                    # Never a static image, so it's always editable.
                    opt_raw = opt.get("raw") or opt.get("display") or ""
                    math_stats = {}
                    # force_align='l' is baked directly into each
                    # paragraph's a:pPr AT XML-BUILD TIME (see
                    # _build_mixed_paragraph_xml) -- this is required
                    # because a paragraph that is ENTIRELY one equation
                    # (a pure-math MCQ option, no leading/trailing text)
                    # can render center-aligned in PowerPoint even when
                    # alignment is only patched on afterwards. Baking it
                    # in up front removes that failure mode.
                    _set_shape_text_with_math(shp, opt_raw, math_stats, force_align='l')
                    _accumulate_math_stats(math_stats)
                    # Belt-and-braces: also set it via python-pptx so any
                    # paragraph missed above (e.g. an empty trailing line)
                    # is still explicitly left-aligned.
                    for para in shp.text_frame.paragraphs:
                        para.alignment = PP_ALIGN.LEFT
                else:
                    # source had fewer than 4 options -- remove the unused
                    # option box + its A/B/C/D circle rather than leaving
                    # a blank "Text here" placeholder visible
                    _remove_shape(shp)
                    grp = _find_shape(dest, Q_OPTION_GROUP_NAMES[i])
                    if grp is not None:
                        _remove_shape(grp)
                    stats["options_short"] += 1
            stats["mcq_slides"] += 1
        else:
            # Numerical / Subjective: no option boxes on this template ->
            # remove all 4 option boxes + their circle badges cleanly.
            for opt_name, grp_name in zip(Q_OPTION_SHAPES, Q_OPTION_GROUP_NAMES):
                shp = _find_shape(dest, opt_name)
                if shp is not None:
                    _remove_shape(shp)
                grp = _find_shape(dest, grp_name)
                if grp is not None:
                    _remove_shape(grp)
            stats["other_q_slides"] += 1

        # If the stem's estimated text height spills down far enough to
        # overlap the option circle-labels/boxes below it, push all the
        # option content (whatever's still left on the slide) down by 2"
        # so the stem never overlaps the circle labels or option text.
        if stem_shape is not None:
            stem_bottom_est = stem_shape.top + max(
                stem_shape.height, _estimate_text_height_emu(stem_shape, stem_raw))
            remaining = [n for n in (Q_OPTION_SHAPES + Q_OPTION_GROUP_NAMES)
                         if _find_shape(dest, n) is not None]
            tops = [_find_shape(dest, n).top for n in remaining]
            if tops and stem_bottom_est > min(tops):
                shift = PPTX_Inches(2)
                for n in remaining:
                    s = _find_shape(dest, n)
                    s.top = s.top + shift

    current_heading = "Notes"
    body_buffer = ""  # accumulates RAW text (math delimiters intact)

    def _split_text_to_fit_theory(text):
        """Split only at safe text boundaries. Native Word OMML tokens are
        atomic and are never cut; visible text is used for height estimation."""
        text=text or ''
        if not text.strip(): return '', ''
        token_store=[]
        def protect(m):
            token_store.append(m.group(0)); return f'OMMLTOKEN{len(token_store)-1}X'
        protected=OMML_TOKEN_RE.sub(protect,text) if 'OMML_TOKEN_RE' in globals() else text
        probe=theory_body_bp_shape
        usable_h=max(1,prs.slide_height-probe.top-int(0.55*914400))
        def restore(s):
            for i,tok in enumerate(token_store): s=s.replace(f'OMMLTOKEN{i}X',tok)
            return s
        if _estimate_text_height_emu(probe,protected)<=usable_h:
            return restore(protected.strip()),''
        words=re.findall(r'\S+|\s+', protected)
        if len(words)<=1:
            lo,hi=1,len(protected); best=1
            while lo<=hi:
                mid=(lo+hi)//2
                if _estimate_text_height_emu(probe,protected[:mid])<=usable_h: best=mid; lo=mid+1
                else: hi=mid-1
            # Never cut through a protected OMML placeholder.
            cut=best
            for i in range(len(token_store)):
                marker=f'OMMLTOKEN{i}X'
                a=protected.find(marker); b=a+len(marker) if a>=0 else -1
                if a<cut<b: cut=a
            if cut<=0: cut=best
            return restore(protected[:cut].rstrip()),restore(protected[cut:].lstrip())
        lo,hi=1,len(words); best=1
        while lo<=hi:
            mid=(lo+hi)//2
            # Tokens alternate word/whitespace and preserve the ORIGINAL
            # separators verbatim (including '\n' paragraph/bullet
            # boundaries) -- unlike '.split()' + ' '.join(...), which
            # silently flattened every newline to a single space and is
            # exactly what used to weld every bullet/paragraph back into
            # one run-on line whenever a slide's content needed to be
            # split across multiple slides.
            candidate=''.join(words[:mid])
            if _estimate_text_height_emu(probe,candidate)<=usable_h: best=mid; lo=mid+1
            else: hi=mid-1
        prefix=''.join(words[:best]).strip(); rest=''.join(words[best:]).strip()
        # Do not split a LaTeX span, or a **bold**/*italic* markdown span,
        # in half -- otherwise the closing marker lands on the NEXT slide
        # and the opening delimiter is left dangling as literal, visible
        # characters instead of rendering as the equation/bold/italic it
        # was meant to be.
        for pattern in (MATH_SPAN_PATTERN, _MD_STYLE_PATTERN):
            combined = prefix + (' ' + rest if rest else '')
            for m in pattern.finditer(combined):
                cutpos = len(prefix)
                if m.start() < cutpos < m.end():
                    safe = combined.rfind(' ', 0, m.start()) if m.start() > 0 else 0
                    if safe > 0:
                        prefix = combined[:safe].strip()
                        rest = combined[safe:].strip()
                    break
        return restore(prefix),restore(rest)


    def flush_body_buffer(force_final=False):
        nonlocal body_buffer
        usable_h = max(1, prs.slide_height - theory_body_bp_shape.top - int(0.55 * 914400))
        # IMPORTANT: normal paragraph/Enter boundaries do not flush a slide.
        # We create a slide only when the accumulated text actually exceeds
        # the available body height, or when a structural block/section ends.
        while body_buffer.strip() and _estimate_text_height_emu(theory_body_bp_shape, body_buffer) > usable_h:
            chunk, remainder = _split_text_to_fit_theory(body_buffer)
            if not chunk:
                break
            new_theory_text_slide(current_heading, chunk)
            body_buffer = remainder
            if not remainder:
                break
        if force_final and body_buffer.strip():
            new_theory_text_slide(current_heading, body_buffer.strip())
            body_buffer = ''

    def take_body_buffer():
        """Pulls out and clears whatever theory text is currently
        buffered, WITHOUT giving it its own slide -- used so a table (or
        other block) that immediately follows can carry that leading text
        onto its own slide instead."""
        nonlocal body_buffer
        text = body_buffer.strip()
        body_buffer = ""
        return text

    for seg in segments:
        if seg["kind"] == "heading":
            flush_body_buffer(force_final=True)
            current_heading = seg["text"] or current_heading
            stats["headings"] += 1
        elif seg["kind"] == "text":
            raw_piece = seg.get("raw", seg["text"])
            # Each "text" segment is already ONE distinct source block
            # (a paragraph or a bullet/list item straight from MinerU's
            # content_list.json, or one Word paragraph on the edited-
            # Word re-upload path). Joining these with a bare space used
            # to weld every bullet point / paragraph into one run-on
            # line once rendered on the slide (bullets showing but with
            # no line breaks between them). Join with a newline instead
            # so _set_shape_text_with_math() gives each one its own
            # paragraph, preserving the source's bullet/paragraph
            # structure exactly as extracted.
            body_buffer = (body_buffer + "\n" + raw_piece).strip() if body_buffer else raw_piece
            flush_body_buffer(force_final=False)
        elif seg["kind"] == "equation":
            flush_body_buffer(force_final=True)
            new_theory_equation_slide(current_heading, seg["latex"])
        elif seg["kind"] == "image":
            flush_body_buffer(force_final=True)
            new_theory_image_slide(current_heading, seg.get("img_path"))
        elif seg["kind"] == "table":
            # Keep the paragraph that leads into this table on the SAME
            # slide as the table itself, instead of flushing it to its
            # own slide first (which pushed the table to the next one).
            leading_text_raw = take_body_buffer()
            new_theory_table_slide(current_heading, seg["rows"], leading_text_raw)
        elif seg["kind"] == "question":
            flush_body_buffer(force_final=True)
            new_question_slide(seg)

    flush_body_buffer(force_final=True)

    # Blueprints were only templates -- remove them from the final deck.
    # Delete the higher index first so the lower index doesn't shift.
    _delete_slide(prs, max(theory_bp_idx, qna_bp_idx))
    _delete_slide(prs, min(theory_bp_idx, qna_bp_idx))

    prs.save(output_pptx_path)
    return stats


def process_pdf(pdf_file, mode_choice, lang_choice):
    if pdf_file is None:
        return None, "Kripya PDF file upload karein.", None, None

    pdf_path = pdf_file.name
    base_name = os.path.splitext(os.path.basename(pdf_path))[0]
    output_workspace = os.path.join(os.getcwd(), "output_workspace")

    if os.path.exists(output_workspace):
        shutil.rmtree(output_workspace)
    os.makedirs(output_workspace, exist_ok=True)

    mode = EXTRACTION_MODES.get(mode_choice, "auto_detect")
    font_risk, risky_fonts = detect_symbol_font_risk(pdf_path)

    # --m ocr: bypasses the PDF's embedded text layer entirely and runs
    #   full layout-detection + visual OCR + formula-recognition on every
    #   page image. This is the ONLY reliable fix when the source PDF's
    #   equation fonts (SymbolMT/MT-Extra/MathType etc.) have a broken
    #   ToUnicode map -- the embedded "text" is unrecoverable garbage no
    #   matter which text-extraction tool reads it, and it also fixes
    #   two-column reading-order issues since layout is re-derived visually.
    # --l <lang>: tunes the OCR model for the document's script.
    if mode == "force_ocr":
        method = "ocr"
    elif mode == "force_txt":
        method = "auto"
    else:  # auto_detect
        method = "ocr" if font_risk else "auto"

    mineru_args = [
        "mineru", "-p", pdf_path, "-o", output_workspace,
        "-b", "pipeline",
        "-m", method,
        "-f", "true",   # formula parsing ON
        "-t", "true",   # table parsing ON
    ]
    lang_code = LANG_MAP.get(lang_choice)
    if lang_code:
        mineru_args += ["-l", lang_code]

    _t_extract_start = time.time()
    result = subprocess.run(mineru_args, capture_output=True, text=True)
    extraction_seconds = time.time() - _t_extract_start

    content_list_files = glob.glob(os.path.join(output_workspace, "**", "*_content_list.json"), recursive=True)
    md_files = glob.glob(os.path.join(output_workspace, "**", "*.md"), recursive=True)
    if not content_list_files and not md_files:
        return None, f"Extraction Failed: {result.stderr[:500]}", None, None

    final_docx_name = f"{base_name}_Enterprise_Converted.docx"
    final_docx_path = os.path.join(output_workspace, final_docx_name)

    used_structured_path = False
    _t_build_start = time.time()
    try:
        if content_list_files:
            # PREFERRED: uses real bbox positions for true left-column-then-
            # right-column ordering, and clean pre-classified LaTeX/HTML for
            # equations/tables -- no markdown-regex guessing involved.
            content_list_path = content_list_files[0]
            extracted_dir = os.path.dirname(content_list_path)
            images_dir = os.path.join(extracted_dir, "images")
            stats = build_docx_from_content_list(content_list_path, extracted_dir, images_dir, final_docx_path)
            used_structured_path = True
        else:
            md_path = md_files[0]
            extracted_dir = os.path.dirname(md_path)
            images_dir = os.path.join(extracted_dir, "images")
            stats = build_enterprise_docx(md_path, images_dir, final_docx_path)
    except Exception as e:
        return None, f"Word Generation Error: {str(e)}", None, None
    build_seconds = time.time() - _t_build_start

    status_lines = []
    status_lines.append(
        f"⏱ Timing -- PDF extraction (MinerU): {extraction_seconds:.1f}s | "
        f"Word document build: {build_seconds:.1f}s | Total: {extraction_seconds + build_seconds:.1f}s"
    )
    if extraction_seconds > 30 and build_seconds < 5:
        status_lines.append(
            "Note: almost all the time above is MinerU's own PDF layout/OCR extraction, "
            "not this app's document-building code (that part finished in a few seconds). "
            "MinerU speed depends on your machine's GPU/CPU and whether its models are "
            "already warmed up -- see the message below this status box for ways to speed "
            "that stage up."
        )
    if font_risk:
        status_lines.append(
            f"Corrupted Symbol/MathType font detect hui ({', '.join(risky_fonts)}) — "
            f"is PDF ka embedded text-layer bharosemand nahi hai, isliye Full OCR mode "
            f"{'auto-force' if mode == 'auto_detect' else 'use'} kiya gaya."
        )
    status_lines.append(f"Extraction method used: {method.upper()}")
    status_lines.append(
        "Content source: content_list.json (STRICT page-by-page geometry order: single-column top-to-bottom; two-column left-to-bottom then right-to-bottom; no content blocks intentionally discarded)"
        if used_structured_path else
        "Content source: markdown (content_list.json not found -- older MinerU version?)"
    )
    if used_structured_path:
        status_lines.append(
            f"Content blocks processed: {stats.get('total_blocks_seen', 0)} "
            f"(every block accounted for -- none silently dropped)"
        )
        if stats.get("crash_recovered"):
            status_lines.append(
                f"⚠ {stats['crash_recovered']} block(s) hit an unexpected error and were "
                f"recovered as flagged plain text -- search the doc for 'manual review needed'"
            )
    if stats.get("skipped_types"):
        status_lines.append(f"Unclassified block types encountered (rendered as plain text, not dropped): {stats['skipped_types']}")
    status_lines.append(
        f"Equations converted to native OMML: {stats['success']} | "
        f"Fallback (shown as LaTeX text, NOT dropped): {stats['failed']}"
    )
    if stats["failed_list"]:
        status_lines.append("Manual-review needed for: " + " | ".join(stats["failed_list"]))

    returned_content_list_path = content_list_files[0] if (used_structured_path and content_list_files) else None
    returned_base_dir = extracted_dir if used_structured_path else None
    return final_docx_path, "\n".join(status_lines), returned_content_list_path, returned_base_dir


def process_pptx_generation(content_list_path, base_dir, template_file, edited_word_file=None):
    """Generate PPTX from the edited Word file when supplied; otherwise keep
    the original content_list.json workflow for backward compatibility."""
    if template_file is None:
        return None, (
            "Kripya apna PPTX template upload karo. Isme 2 slides honi chahiye:\n"
            "1st slide (Theory blueprint) -- shapes named 'Rectangle: Rounded Corners 1' (heading) aur 'TextBox 1' (body)\n"
            "2nd slide (Q&A blueprint) -- shapes named 'TextBox 21'/'22'/'23'/'24'/'25'"
        )

    output_workspace = os.path.join(os.getcwd(), "output_workspace")
    os.makedirs(output_workspace, exist_ok=True)
    output_pptx_path = os.path.join(output_workspace, "Generated_Presentation.pptx")

    try:
        _t0 = time.time()
        if edited_word_file is not None:
            word_path = edited_word_file.name
            word_workspace = os.path.join(output_workspace, 'edited_word_workspace')
            if os.path.exists(word_workspace):
                shutil.rmtree(word_workspace)
            os.makedirs(word_workspace, exist_ok=True)
            stats = build_pptx_from_word(word_path, template_file.name, output_pptx_path, word_workspace)
            source_label = "EDITED WORD (.docx)"
        else:
            if not content_list_path:
                return None, (
                    "Pehle PDF ko Word me convert karo, ya directly 'Edited Word File' upload karo. "
                    "Edited Word upload karne par wahi file PPTX ka source banegi."
                )
            stats = build_pptx_from_content_list(content_list_path, base_dir, template_file.name, output_pptx_path)
            source_label = "content_list.json (legacy/backward-compatible mode)"
        pptx_seconds = time.time() - _t0
    except Exception as e:
        return None, f"PPTX Generation Error: {type(e).__name__}: {str(e)}"

    status_lines = [
        f"Source: {source_label}",
        f"⏱ PPTX build time: {pptx_seconds:.1f}s",
        f"Theory slides generated: {stats['theory_slides']} (across {stats['headings']} detected heading(s))",
        f"MCQ slides generated: {stats['mcq_slides']}",
        f"Numerical/Subjective slides generated: {stats['other_q_slides']}",
        f"Equations rendered as REAL editable PowerPoint equations (OMML): {stats['equations_as_omml']}",
        f"Images/diagrams embedded: {stats['images_embedded']}",
        f"Tables embedded as real PowerPoint tables: {stats['tables_embedded']}",
        "Heading geometry/formatting is cloned from the PPTX template; Word Enter does NOT force a new slide.",
        "Theory content moves to the next slide only when the actual available TextBox 1 height is exceeded.",
    ]
    if stats.get("equations_as_text_fallback"):
        status_lines.append(f"Equation text fallbacks: {stats['equations_as_text_fallback']}")
    if stats.get("images_missing"):
        status_lines.append(f"⚠ Missing images: {stats['images_missing']}")
    if stats.get("options_short"):
        status_lines.append(f"⚠ Short-option questions: {stats['options_short']}")
    return output_pptx_path, "\n".join(status_lines)


# Gradio Interface
with gr.Blocks(title="PDF to Word/PPTX Enterprise Native Engine") as app:
    gr.Markdown("## PDF to Word / PowerPoint Enterprise Native Engine")
    gr.Markdown(
        "Native Python OMML transpiler (Vector / Sum / Integral / Array equations), "
        "auto-detection of corrupted MathType/Symbol fonts, self-calibrating layout-aware "
        "reading order (works on single-column, multi-column, and scanned PDFs alike), "
        "source-accurate MCQ numbering, and template-driven PPTX generation with real "
        "embedded equations/images/tables."
    )
    content_list_state = gr.State(None)
    base_dir_state = gr.State(None)

    with gr.Tab("1. PDF to Word"):
        with gr.Row():
            pdf_input = gr.File(label="Upload PDF File (Vectors / Arrays / Math / MCQs)", file_types=[".pdf"])
        with gr.Row():
            mode_dropdown = gr.Dropdown(
                choices=list(EXTRACTION_MODES.keys()),
                value="Auto-detect (recommended)",
                label="Extraction Mode"
            )
            lang_dropdown = gr.Dropdown(
                choices=list(LANG_MAP.keys()), value="Auto / English (default)",
                label="Document Language (OCR accuracy tuning)"
            )
        convert_btn = gr.Button("Convert to Word", variant="primary")
        output_file = gr.File(label="Download Enterprise Word (.docx) File")
        output_status = gr.Textbox(label="Status", lines=5)

        convert_btn.click(
            fn=process_pdf,
            inputs=[pdf_input, mode_dropdown, lang_dropdown],
            outputs=[output_file, output_status, content_list_state, base_dir_state]
        )

    with gr.Tab("2. Generate PPTX from Template"):
        gr.Markdown(
            "Pehle Tab 1 me ek PDF convert karo, phir yahan apna PPTX template upload karo. "
            "Template me 2 slides honi chahiye: **Slide 1 (Theory)** with shapes named "
            "`Rectangle: Rounded Corners 1` (heading) + `TextBox 1` (body), and **Slide 2 (Q&A)** "
            "with shapes named `TextBox 21` (question) + `TextBox 22/23/24/25` (options A-D). "
            "These two slides repeat automatically for the whole document -- theory pages get "
            "auto-continued onto new slides when content overflows, and every MCQ/Numerical/"
            "Subjective question gets exactly one slide. Equations render as real math images, "
            "diagrams/tables are embedded for real -- nothing is converted to a plain-text note."
        )
        template_input = gr.File(label="Upload PPTX Template", file_types=[".pptx"])
        edited_word_input = gr.File(
            label="Upload Edited Word File (OPTIONAL - preferred PPTX source)",
            file_types=[".docx"]
        )
        gr.Markdown(
            "**PPTX source priority:** Edited Word file > original `content_list.json`. "
            "Agar aapne generated Word file me heading, paragraph, Enter/line-break, image, "
            "table ya text correction ki hai, to yahi edited `.docx` upload karein."
        )
        generate_pptx_btn = gr.Button("Generate PPTX", variant="primary")
        pptx_output_file = gr.File(label="Download Generated PPTX")
        pptx_output_status = gr.Textbox(label="Status", lines=8)

        generate_pptx_btn.click(
            fn=process_pptx_generation,
            inputs=[content_list_state, base_dir_state, template_input, edited_word_input],
            outputs=[pptx_output_file, pptx_output_status]
        )

if __name__ == "__main__":
    app.launch(inbrowser=True)
